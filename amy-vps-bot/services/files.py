from __future__ import annotations

import io
from html import escape

import docx
import pdfplumber
from pptx import Presentation
from pptx.util import Inches, Pt
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER
from reportlab.lib.pagesizes import A4
from reportlab.lib.styles import ParagraphStyle
from reportlab.lib.units import cm
from reportlab.pdfbase import pdfmetrics
from reportlab.pdfbase.cidfonts import UnicodeCIDFont
from reportlab.platypus import HRFlowable, Paragraph, SimpleDocTemplate, Spacer


def extract_text(file_bytes: bytes, mime: str, filename: str) -> str:
    if mime == "application/pdf" or filename.lower().endswith(".pdf"):
        with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
            return "\n".join(page.extract_text() or "" for page in pdf.pages)[:20000]

    if mime.startswith("text/") or filename.lower().endswith((".txt", ".csv", ".md", ".json")):
        return file_bytes.decode("utf-8", errors="ignore")[:20000]

    if filename.lower().endswith(".docx"):
        doc = docx.Document(io.BytesIO(file_bytes))
        return "\n".join(p.text for p in doc.paragraphs if p.text.strip())[:20000]

    return ""


def make_pdf(title: str, content: str) -> bytes:
    pdfmetrics.registerFont(UnicodeCIDFont("STSong-Light"))
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(
        buffer,
        pagesize=A4,
        rightMargin=2.2 * cm,
        leftMargin=2.2 * cm,
        topMargin=2.2 * cm,
        bottomMargin=2.2 * cm,
    )
    title_style = ParagraphStyle(
        "Title",
        fontName="STSong-Light",
        fontSize=18,
        leading=24,
        alignment=TA_CENTER,
        spaceAfter=12,
        textColor=colors.HexColor("#111827"),
    )
    body_style = ParagraphStyle("Body", fontName="STSong-Light", fontSize=11, leading=18, spaceAfter=6)
    heading_style = ParagraphStyle(
        "Heading",
        fontName="STSong-Light",
        fontSize=13,
        leading=20,
        spaceBefore=12,
        spaceAfter=6,
        textColor=colors.HexColor("#1d4ed8"),
    )

    story = [Paragraph(escape(title), title_style), HRFlowable(width="100%", color=colors.lightgrey), Spacer(1, 12)]
    for line in content.splitlines():
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 6))
        elif stripped.startswith("#") or stripped.endswith("："):
            story.append(Paragraph(escape(stripped.lstrip("#").strip()), heading_style))
        else:
            story.append(Paragraph(escape(stripped), body_style))
    doc.build(story)
    return buffer.getvalue()


def make_pptx(title: str, slides: list[dict[str, list[str] | str]]) -> bytes:
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    first = prs.slides.add_slide(prs.slide_layouts[0])
    first.shapes.title.text = title
    first.shapes.title.text_frame.paragraphs[0].font.size = Pt(38)

    for item in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = str(item["title"])
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for point in item["points"]:
            p = tf.add_paragraph()
            p.text = str(point)
            p.font.size = Pt(21)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()
