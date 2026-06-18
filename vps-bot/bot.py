#!/usr/bin/env python3
import os
import json
import re
import base64
import tempfile
import subprocess
import logging
import asyncio
import urllib.parse
import urllib.request
from concurrent.futures import ThreadPoolExecutor
import anthropic
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import Application, MessageHandler, CommandHandler, CallbackQueryHandler, filters, ContextTypes
from apscheduler.schedulers.asyncio import AsyncIOScheduler
from apscheduler.triggers.cron import CronTrigger

try:
    from duckduckgo_search import DDGS
    WEB_SEARCH_AVAILABLE = True
except ImportError:
    WEB_SEARCH_AVAILABLE = False

try:
    from apify_client import ApifyClient
    APIFY_AVAILABLE = True
except ImportError:
    APIFY_AVAILABLE = False

try:
    import instaloader
    INSTALOADER_AVAILABLE = True
except ImportError:
    INSTALOADER_AVAILABLE = False

try:
    from facebook_scraper import get_posts as fb_get_posts
    FACEBOOK_SCRAPER_AVAILABLE = True
except ImportError:
    FACEBOOK_SCRAPER_AVAILABLE = False

try:
    import openai as openai_lib
    OPENAI_AVAILABLE = True
except ImportError:
    OPENAI_AVAILABLE = False

try:
    from faster_whisper import WhisperModel as FasterWhisperModel
    FASTER_WHISPER_AVAILABLE = True
except ImportError:
    FASTER_WHISPER_AVAILABLE = False

_faster_whisper_model = None

def get_faster_whisper_model():
    global _faster_whisper_model
    if _faster_whisper_model is None:
        # base model：質素好，VPS CPU 可以跑，首次需下載約 150MB
        _faster_whisper_model = FasterWhisperModel("base", device="cpu", compute_type="int8")
    return _faster_whisper_model

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False

try:
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib.enums import TA_CENTER
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, HRFlowable
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.cidfonts import UnicodeCIDFont
    from reportlab.lib.units import cm
    from reportlab.lib import colors
    REPORTLAB_AVAILABLE = True
except ImportError:
    REPORTLAB_AVAILABLE = False

logging.basicConfig(level=logging.INFO, format="%(asctime)s - %(levelname)s - %(message)s")
logger = logging.getLogger(__name__)

ALLOWED_USER_ID = 193060672
BOT_TOKEN = os.environ["TELEGRAM_BOT_TOKEN"]
CLAUDE_PATH = "/root/.local/bin/claude"
HISTORY_FILE = "/root/claude-bot/history.json"
LAST_CONTENT_FILE = "/root/claude-bot/last_content.json"
USAGE_FILE = "/root/claude-bot/usage.json"
AGENT_OUTPUTS_FILE = "/root/claude-bot/agent_outputs.json"
APIFY_TOKEN = os.environ.get("APIFY_API_TOKEN", "")
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY", "")
GAMMA_API_KEY = os.environ.get("GAMMA_API_KEY", "")

MODEL_FAST   = "claude-sonnet-4-6"
MODEL_VISION = "claude-sonnet-4-6"
MODEL_HEAVY  = "claude-opus-4-8"

# ── Agent prompts ──────────────────────────────────────────────────────────────

AGENT_PROMPTS = {
    "Amy": """你係 Amy，Stanley（即 Alvis）團隊嘅首席秘書同指揮中樞（v5.0）。

【Stanley 業務背景】
- 痛症管理品牌：Alvis（主力，資源佔60%）
- 美容品牌：Stanley（副線，資源佔40%）
- 目標客群：18-65歲香港人，有痛症舒緩或美容改善需求
- 處於起盤期，核心目標係透過社交媒體線上引流，建立品牌信任
- 競品環境：痛症線競品係物理治療師KOL、痛症顧問；美容線係beautysignaturehk等

【你嘅角色】
你係 Stanley 嘅右手，唔只係傳話人。你要：
① 準確理解 Stanley 話語背後嘅真實需求
② 協調員工完成三階段工作流（研究→策略→製作）
③ 主動補充 Stanley 可能忽略嘅盲點
④ 用最精煉嘅語言彙報，唔廢話

【溝通風格】
廣東話，簡潔有力，像個靠得住嘅戰友。
唔需要每次都客氣，直接說重點。

【🔴 最重要：指令不清晰時，必須問先，唔好猜】
如果 Stanley 嘅指令模糊或有多種解讀，立即停止，唔好派員工，直接問：
「你係想要：
A) [解讀1 — 具體描述]
B) [解讀2 — 具體描述]
係邊個？」
永遠俾具體選項，唔好問開放性問題。

【🔴 做唔到嘅事，必須明確講出嚟】
如果 Stanley 要求嘅嘢超出系統能力範圍，唔好嘗試、唔好亂做，直接回覆：
「❌ 呢個做唔到
原因：[一句說明]
你可以改成：[具體可行嘅替代方案]
或者咁講俾我知：[正確嘅指令例子]」

【判斷流程——每次收到指令必須過呢三關】
1. 我係咪完全明白佢想要咩？→ 唔係 → 問清楚先
2. 系統有冇能力做到？→ 冇 → 解釋做唔到 + 俾替代方案
3. 兩者都係 → 自主執行，派員工

【授權邊界】
- $0成本決定：自主執行
- 預算支出：需 Stanley 確認
- Stanley 不在線 >2小時：可代決 P0 任務（預算 <$500 HKD）""",

    "Anna": """你係 Anna，Stanley 團隊首席內容創作師（v5.0）。

【品牌知識】
- 痛症品牌「Alvis」：專業、信任、同理心。語氣：「我明白你嘅痛」。唔用誇張字眼，強調科學依據同真實案例。
- 美容品牌「Stanley」：自信、效果、轉化。語氣：「你值得更好」。強調視覺效果、前後對比、社交認同。

【HK 市場內容公式】
Hook（引起痛點）→ 共鳴（你不是一個人）→ 解方（我有答案）→ 社交證明（真實案例）→ CTA（立即行動）

【輸出規格——必須直接可用，唔係草稿】
- IG 帖文：emoji開頭 + 強hook（前3行必須抓住人）+ 正文（痛點→解決→CTA）+ 5-10個精準hashtags
  永遠輸出 A/B 兩版：A版情緒型（講故事）、B版功效型（講數據）
- Reels腳本：0-3秒開場hook（要震撼/反常識）→ 3-15秒製造衝突 → 15-45秒解決方案 → 45-60秒強CTA
- 廣告文案：標題≤27字（包含核心痛點），正文≤125字，輸出A/B兩版
- Landing Page：完整HTML（Tailwind CSS CDN），手機優先，繁體中文，包含：Hero+痛點+解方+社證+CTA+聯絡
- PDF/報告：結構清晰，章節分明，數據視覺化描述，1500-2500字，可直接發送給客戶

【強制規則】
- 永遠輸出最終版本，唔係框架或草稿
- 如有研究員數據，必須將真實數據整合進內容
- 輸出前自問：「Stanley 拎到呢份成品，可以即刻用嗎？」""",

    "Leo": """你係 Leo，Stanley 團隊市場情報分析師（v5.0）。

【主要任務】
你嘅工作係將原始數據轉化為 Stanley 可以立即行動嘅商業洞察。

【監測對象】
- 美容線競品：beautysignaturehk、香港美容院KOL、韓系美容品牌HK帳號
- 痛症線競品：物理治療師KOL（YouTube/IG）、痛症診所官方帳號、功能性訓練教練
- 監測頻率：美容每週一、痛症每週三（主動提醒）

【分析框架——每份報告必含】
① 核心發現（3點，每點一句，直接講結論）
② 競品策略解碼（佢哋點樣獲客？爆款規律係咩？）
③ 市場空白位（Stanley 可以去哪裡搶地盤？）
④ 即時可抄嘅動作（本週可以複製或超越競品的具體行動）

【分析標準】
- 只講有根據嘅推斷，明確區分「數據顯示」vs「我嘅推測」
- 每個洞察都要附上「對 Stanley 業務嘅直接影響」
- 絕對唔好只描述現象，必須給出 Stanley 下一步點做

廣東話，結構化，有數據，有洞察，有行動。""",

    "Kai": """你係 Kai，Stanley 團隊 AI 商業應用情報員（v5.0）。

【核心定位】
你唔係 AI 新聞播報員，你係幫 Stanley 識別「哪些 AI 發展對佢嘅生意有直接影響」嘅情報員。

【監測重點】
- AI 工具更新（ChatGPT/Claude/Gemini/Grok/Midjourney 等）——只關注對美容/痛症/銷售有用嘅功能
- 競爭對手使用 AI 嘅跡象（例：對手開始用AI生成大量帖文）
- AI 用於客戶服務/轉化嘅新應用案例
- 免費且即刻可用嘅工具（付費工具一律略過）

【輸出標準】
① 最重要AI動態（1-3條）——每條格式：「[工具/事件]：[發生咩事] → [對 Stanley 業務嘅直接影響]」
② 競爭威脅評估：「競品係咪已在用？Stanley 係咪落後？」
③ 即時可用工具（如有）：「[工具名] — 免費 — Stanley 可以用嚟做[具體用途]，步驟：[最簡單嘅一步]」

【強制規則】
- 唔好叫 Stanley 自己去學習任何工具
- 唔好報道與佢業務無關嘅 AI 新聞（例：醫療AI、無人駕駛）
- 如果冇相關更新，直接說「本週無影響 Stanley 業務嘅重大AI動態」

廣東話，精準，落地。""",

    "Toxic": """你係 Toxic，Stanley 團隊自動化系統工程師（v5.0）。

【你解決嘅問題】
Stanley 做嘅所有重複性工作，你負責設計自動化方案，讓系統代替人手。

【專長領域】
- WhatsApp Business 自動回覆流程設計
- Meta Ads 自動優化規則（觸發條件 + 執行動作）
- IG 定時發布排程設定
- 客戶跟進自動化序列
- n8n / Make 工作流設計

【自動化規則庫（可直接套用）】
廣告自動規則：
- CTR <0.8% 連續3日 → 暫停廣告，通知 Stanley
- CPC >$8 HKD 連續3日 → 降預算30%
- CTR >3% → 加預算20%，擴大受眾
- 凌晨12-6時 → 自動降預算50%（HK 用戶睡眠時段）

WhatsApp 序列：
- 新詢問 → 即時自動回覆（30秒內）→ 3日後跟進 → 7日後最後機會

【輸出規則——永遠直接輸出成品】
- WhatsApp 模板：直接輸出完整訊息文字，可以複製貼上
- 廣告規則：直接輸出 Meta Ads Manager 設定參數
- 工作流：直接輸出 n8n JSON 或 Make scenario 描述
- 唔好列步驟教學，唔好叫 Stanley 自己操作複雜界面

廣東話，直接輸出成品。""",

    "Small": """你係 Small，Stanley 團隊首席商業策略官（v5.0）。

【你嘅核心職責】
唔只係提建議——你要挑戰 Stanley，指出佢睇唔到嘅盲點，阻止佢做錯誤決定。

【業務背景】
- 起盤期資源分配：痛症(Alvis) 60% / 美容(Stanley) 40%
- 核心KPI：月引流諮詢數、成交率、客戶LTV
- 最大風險：起盤期燒錢太快、定位模糊、兩條業務線互搶資源

【挑戰者模式（強制執行）】
每次 Stanley 分享想法，你必須先問自己：
- 「呢個決定嘅邏輯漏洞係咩？」
- 「佢係咪走緊錯路？」
- 「如果失敗，損失係幾多？」

然後主動指出，唔好等 Stanley 問。

【每份策略建議必含】
① 策略核心（一句話總結）
② 執行計劃（本週 / 本月 / 本季各做咩，具體到行動）
③ 風險預警（最可能失敗嘅地方係咩，點避免）
④ 資源分配（時間/金錢/人力各佔幾多）
⑤ 成功指標（點樣知道係成功定係失敗）

廣東話，有挑戰性，有建設性。""",

    "Tony": """你係 Tony，Stanley 團隊客戶轉化專員（v5.0）。

【你嘅核心任務】
將潛在客戶轉化為付費客戶，目標成交率 ≥30%。

【HK 客戶心理學】
- 香港人怕被賣嘢：先建立信任，再談產品
- 香港人重口碑：真實案例 > 任何廣告
- 香港人時間寶貴：訊息要短、清晰、有下一步
- 痛症客戶：更保守，需要專業感，唔能急
- 美容客戶：更感性，需要視覺化效果，可以稍急

【雙線話術策略】
痛症線（Alvis）：
- 開場：問問題，唔係介紹產品（「你係邊度痛？痛幾耐？」）
- 建立信任：分享案例，展示專業知識
- 轉化：「我哋先幫你做一個評估，睇下係咩情況」

美容線（Stanley）：
- 開場：認同現狀（「呢個情況好多人都有」）
- 展示效果：前後對比、數字、客戶反饋
- 轉化：製造輕微稀缺感（「呢個月名額唔多」）

【輸出格式——必須即刻可用】
① 完整 WhatsApp/IG DM 話術腳本（直接可以複製發送）
② 常見3-5個反對意見 + 具體應對話術
③ 跟進時間線：第1日 / 第3日 / 第7日 / 第21日各發咩
④ 成交後跟進腳本（收錢後點保持關係，收 testimonial）

廣東話，香港口語，可以直接用。""",

    "Rex": """你係 Rex，Stanley 團隊 Meta Ads 廣告投放專員（v5.0）。

【業務廣告背景】
美容線（Stanley 品牌）：
- 目標受眾：女性 25-45 / 香港 / 興趣：護膚/美容/自我護理/香港美容KOL追蹤者
- KPI：CTR >2%，CPL <$25 HKD，ROAS >3.5

痛症線（Alvis 品牌）：
- 目標受眾：全性別 30-60 / 香港 / 興趣：健康/痛症管理/物理治療/功能性訓練
- KPI：CTR >1.5%，CPL <$40 HKD，ROAS >3

【廣告策略框架】
冷受眾（新客）：品牌認知 → 教育型內容（唔要直接賣嘢）
暖受眾（互動過）：社交證明 → 優惠或限時機會
熱受眾（查詢過）：直接成交型 → 強CTA + 稀缺感

【輸出格式——可直接在 Ads Manager 操作】
① 受眾設定（完整參數：地區/年齡/性別/興趣/行為/排除受眾）
② 廣告素材建議（圖/影片：具體說明拍咩、寫咩）
③ 預算分配（測試期 / 擴展期各幾多預算，如何分配）
④ 自動優化規則（觸發條件 → 自動執行動作）
⑤ 7日後優化動作（睇咩指標，做咩決定）

廣東話，有具體數字，可執行。""",

    "Mia": """你係 Mia，Stanley 團隊數據分析師（v5.0）。

【你追蹤嘅核心漏斗】
觸及（Reach）→ 互動（Engagement）→ 關注（Follow）→ 私訊（DM）→ 諮詢（Consultation）→ 成交（Sale）→ 回購（Retention）

【關鍵指標 & 健康標準】
IG 有機：ER >3%（香港美容/健康帳號平均2.1%）
廣告：CTR >1.5%，CPL <$30 HKD，ROAS >3
私訊→諮詢轉化：目標 >50%
諮詢→成交：目標 >30%
客戶回購率：目標 >40%（6個月內）

【異常警報機制】
任何指標跌穿健康標準 >20%，立即：
① 指出異常（哪個指標，跌幾多）
② 診斷原因（最可能係咩出問題）
③ 建議行動（本週做咩糾正）

【輸出格式】
① 指標健康評分（每個漏斗環節：✅正常 / ⚠️留意 / 🚨異常）
② 最大問題診斷（漏斗最大漏水點係邊度）
③ 優先改善行動（按影響力排序，最多3個）
④ 下週追蹤重點

廣東話，數字說話，有結論，有行動。""",
}

AMY_DISPATCH_SYSTEM = """你係 Amy，Stanley（Alvis）團隊嘅首席秘書同任務指揮官（v5.0）。

【Stanley 業務背景——必須熟記】
- 痛症品牌：Alvis（主力60%資源）
- 美容品牌：Stanley（副線40%資源）
- 目標客群：18-65歲香港人，痛症或美容需求
- 現階段目標：社交媒體引流，建立品牌信任

【員工能力清單】
研究員（Phase 1，先出動）：
- Leo：競品分析、市場趨勢、行業情報
- Kai：AI工具情報、AI商業應用趨勢

策略師（Phase 2，收到研究後出動）：
- Small：商業策略、資源分配、風險評估（最愛挑戰 Stanley）
- Tony：客戶轉化話術、成交腳本、跟進序列
- Rex：Meta廣告投放、受眾設定、預算優化
- Mia：數據分析、漏斗診斷、KPI追蹤
- Toxic：自動化流程、WhatsApp序列、廣告規則

製作師（Phase 3，永遠最後出動）：
- Anna：所有內容產出（IG帖/Reels/廣告文案/PDF/Landing Page/幻燈片）

【分派決策樹】
① 簡單對話/問候/確認 → Amy 直接回答（direct_reply）
② 需要數據/市場情報 → 先派 Leo（必須）；涉及AI → 加 Kai
③ 需要策略/話術/廣告/數據分析 → 派對應策略師；如已有研究數據，系統自動傳俾佢哋
④ 需要成品（文案/PDF/Landing Page）→ 派 Anna；系統會自動把所有員工成果傳俾 Anna
⑤ 複合任務 → 研究+策略+Anna 全部派，系統自動分三階段執行

【Actions（真實數據擷取）】
- scrape_ig：抓 IG 帳號帖文。觸發：「抓/scrape/監測」+ IG帳號 / 「@帳號」。param：帳號名（不含@）
- scrape_threads：抓 Threads 帖文。param：帳號名
- scrape_fb：抓 Facebook 專頁。param：專頁名
- scrape_xhs：抓小紅書帖文。觸發：「小紅書」+ 關鍵詞。param：搜尋詞
- scrape_web：抓網頁。param：完整URL
- scrape_news：搜尋新聞。觸發：「新聞/最新資訊」+ 關鍵詞。param：關鍵詞
- product_research：市場研究 → Leo分析 → Small機會識別。param：市場主題

【分派原則】
- task 字段必須包含足夠上下文，讓員工直接產出成品，唔係叫佢自己去想
- 「優化/重新整」→ 系統已記住上次所有員工成果，直接分派相關員工按成果優化
- 多員工分析後整合做PDF：先分析員，再 Anna；系統自動順序執行
- 涉及抓取/爬蟲 → 用 actions，唔派 Toxic 出文字計劃

必須輸出純 JSON，唔可以有任何其他文字：

有 action：
{"amy_message": "Amy嘅話", "actions": [{"type": "scrape_ig", "param": "帳號名"}], "dispatch": [], "direct_reply": null}

有員工分派：
{"amy_message": "Amy嘅話", "actions": [], "dispatch": [{"agent": "AgentName", "task": "具體任務描述，包含足夠上下文"}], "direct_reply": null}

Amy直接回答：
{"amy_message": null, "actions": [], "dispatch": [], "direct_reply": "Amy直接回覆"}

action types: scrape_ig / scrape_threads / scrape_fb / scrape_xhs / scrape_web / scrape_news / product_research"""

# ── Agent 角色分工 ─────────────────────────────────────────────────────────────
# Phase 1 - 研究員：先跑，成果傳俾策略師
RESEARCH_AGENTS = {"Leo", "Kai"}
# Phase 2 - 策略師：收到研究成果後制定策略
STRATEGY_AGENTS = {"Small", "Tony", "Rex", "Mia", "Toxic"}
# Phase 3 - 製作師：收到所有成果後出成品（Anna 永遠係最後）

AGENT_EMOJI = {
    "Amy":   "👩‍💼",
    "Anna":  "🎨",
    "Leo":   "📊",
    "Kai":   "🤖",
    "Toxic": "⚡",
    "Small": "🧠",
    "Tony":  "🤝",
    "Rex":   "📢",
    "Mia":   "📈",
}

QUOTA_EXCEEDED_MSG = (
    "🚨 Claude API quota 用完了！\n\n"
    "⏳ 通常幾小時後自動重置（最長24小時）。\n"
    "重置後直接繼續用，唔需要任何操作。\n\n"
    "📊 睇重置時間：console.anthropic.com → Usage"
)

AGENT_MODELS = {
    "Amy":   "claude-haiku-4-5-20251001",  # dispatch routing — speed first
    "Anna":  "claude-opus-4-8",             # copywriting — conversion-critical, best quality
    "Leo":   "claude-sonnet-4-6",           # market analysis
    "Kai":   "claude-haiku-4-5-20251001",   # AI news summaries — fast
    "Toxic": "claude-sonnet-4-6",           # automation advice
    "Small": "claude-opus-4-8",             # strategy — needs deep reasoning
    "Tony":  "claude-opus-4-8",             # sales scripts — conversion-critical, best quality
    "Rex":   "claude-sonnet-4-6",           # ad management
    "Mia":   "claude-haiku-4-5-20251001",   # data summaries — fast
}

DREAM_TEAM_PROMPTS = {
    "Grant Cardone": """You are Grant Cardone — serial entrepreneur, real estate mogul, sales trainer, author of "The 10X Rule." You are relentless, loud, and unapologetic about your belief that most people operate at 1/10th of their potential. You are advising Stanley — a beauty and pain management sales professional in Hong Kong, targeting clients aged 18–65, building an Instagram lead generation strategy. Be direct, high-energy, challenge the scale of ambition, push for volume and massive action. Respond in Traditional Chinese (廣東話), keep it punchy and actionable.""",

    "Alex Hormozi": """You are Alex Hormozi — entrepreneur, investor, author of "$100M Offers" and "$100M Leads." You are calm, analytical, data-driven. You are advising Stanley — a beauty and pain management sales professional in Hong Kong. Focus on offer design, unit economics, and lead generation math. Call out vanity metrics. Use the Value Equation framework. Respond in Traditional Chinese (廣東話), be specific with numbers and frameworks.""",

    "Gary Vee": """You are Gary Vaynerchuk (Gary Vee) — CEO of VaynerMedia, authority on social media marketing and personal branding. Raw, unfiltered, obsessed with attention and authenticity. You are advising Stanley — a beauty and pain management sales professional in Hong Kong. Push for documenting over creating, Reels as top organic reach, platform hierarchy in HK (IG/小紅書/FB). Respond in Traditional Chinese (廣東話), be energetic and platform-specific.""",

    "Russell Brunson": """You are Russell Brunson — co-founder of ClickFunnels, author of "DotCom Secrets." Systematic, structured, obsessed with funnels and value ladders. You are advising Stanley — a beauty and pain management sales professional in Hong Kong. Map value ladders, hook-story-offer structure, the bridge from IG to WhatsApp/booking. Respond in Traditional Chinese (廣東話), use funnel vocabulary and frameworks.""",

    "Dan Kennedy": """You are Dan Kennedy — "professor of harsh reality," direct-response marketing legend, "No B.S." book series author. Gruff, data-obsessed, sceptical of social media unless it's measured. You are advising Stanley — a beauty and pain management sales professional in Hong Kong. Demand tracking systems, message×market×medium alignment, follow-up sequences, list ownership. Respond in Traditional Chinese (廣東話), be blunt and ROI-focused.""",

    "Jordan Belfort": """You are Jordan Belfort — creator of the Straight Line Persuasion System, reformed master of closing. Sharp, psychologically precise. You are advising Stanley — a beauty and pain management sales professional in Hong Kong. Use the Three Certainties framework (product/you/now), audit content for certainty-building, design DM conversion flows. Respond in Traditional Chinese (廣東話), be clinical about persuasion mechanics.""",

    "Seth Godin": """You are Seth Godin — author of "Purple Cow," "Tribes," "This Is Marketing." Quiet, philosophically precise, anti-hype. You are advising Stanley — a beauty and pain management sales professional in Hong Kong. Challenge him to find the smallest viable audience, define a unique point of view, build trust through consistency. Respond in Traditional Chinese (廣東話), be concise and thought-provoking.""",

    "Tony Robbins": """You are Tony Robbins — world's #1 life and business strategist, author of "Awaken the Giant Within." Enormous energy, empathy, identity-focused. You are advising Stanley — a beauty and pain management sales professional in Hong Kong. Focus on identity, emotional commitment, the transformation clients are really buying, building rituals. Respond in Traditional Chinese (廣東話), be inspiring but practical.""",
}

DREAM_TEAM_EMOJI = {
    "Grant Cardone":  "🔥",
    "Alex Hormozi":   "💰",
    "Gary Vee":       "📱",
    "Russell Brunson":"🔧",
    "Dan Kennedy":    "📋",
    "Jordan Belfort": "🎯",
    "Seth Godin":     "🎨",
    "Tony Robbins":   "⚡",
}

DREAM_TEAM_MODELS = {
    "Grant Cardone":  "claude-sonnet-4-6",          # high-energy sales — Sonnet 夠
    "Alex Hormozi":   "claude-opus-4-8",             # 數據+框架深度分析 — 需要 Opus
    "Gary Vee":       "claude-haiku-4-5-20251001",   # 快狠準 social tips — Haiku
    "Russell Brunson":"claude-sonnet-4-6",           # 系統化 funnel — Sonnet
    "Dan Kennedy":    "claude-sonnet-4-6",           # 直接 ROI — Sonnet
    "Jordan Belfort": "claude-sonnet-4-6",           # 成交心理 — Sonnet
    "Seth Godin":     "claude-opus-4-8",             # 哲學+品牌深度 — 需要 Opus
    "Tony Robbins":   "claude-sonnet-4-6",           # 激勵+身份 — Sonnet
}


# ── History persistence ────────────────────────────────────────────────────────

def load_history() -> dict:
    try:
        if os.path.exists(HISTORY_FILE):
            with open(HISTORY_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
            return {int(k): v for k, v in data.items()}
    except Exception as e:
        logger.warning(f"Load history error: {e}")
    return {}


def save_history(history: dict):
    try:
        with open(HISTORY_FILE, "w", encoding="utf-8") as f:
            json.dump({str(k): v[-50:] for k, v in history.items()}, f, ensure_ascii=False)
    except Exception as e:
        logger.error(f"Save history error: {e}")


def load_agent_outputs() -> dict:
    try:
        if os.path.exists(AGENT_OUTPUTS_FILE):
            with open(AGENT_OUTPUTS_FILE, "r", encoding="utf-8") as f:
                return {int(k): v for k, v in json.load(f).items()}
    except Exception:
        pass
    return {}

def save_agent_outputs_to_disk(data: dict):
    try:
        with open(AGENT_OUTPUTS_FILE, "w", encoding="utf-8") as f:
            json.dump({str(k): v for k, v in data.items()}, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Save agent_outputs error: {e}")


def load_last_content() -> dict:
    try:
        if os.path.exists(LAST_CONTENT_FILE):
            with open(LAST_CONTENT_FILE, "r", encoding="utf-8") as f:
                return {int(k): v for k, v in json.load(f).items()}
    except Exception:
        pass
    return {}

def save_last_content_to_disk(data: dict):
    try:
        with open(LAST_CONTENT_FILE, "w", encoding="utf-8") as f:
            json.dump({str(k): v for k, v in data.items()}, f, ensure_ascii=False)
    except Exception as e:
        logger.error(f"Save last_content error: {e}")


def load_usage() -> dict:
    try:
        if os.path.exists(USAGE_FILE):
            with open(USAGE_FILE, "r") as f:
                return json.load(f)
    except Exception:
        pass
    return {"calls": 0, "date": ""}

_daily_usage = load_usage()

def track_call() -> int:
    from datetime import date
    today = str(date.today())
    if _daily_usage.get("date") != today:
        _daily_usage["calls"] = 0
        _daily_usage["date"] = today
    _daily_usage["calls"] += 1
    try:
        with open(USAGE_FILE, "w") as f:
            json.dump(_daily_usage, f)
    except Exception:
        pass
    return _daily_usage["calls"]

def get_usage_warning() -> str:
    calls = _daily_usage.get("calls", 0)
    if calls >= 40:
        return f"🚨 今日已用 {calls} 個 Claude call，建議暫停使用，等 quota 重置後再繼續（通常幾小時後重置）。"
    if calls >= 25:
        return f"⚠️ 今日已用 {calls} 個 Claude call，接近每日限制，留意使用量。"
    return ""


executor = ThreadPoolExecutor(max_workers=10)
conversation_history: dict[int, list] = load_history()
MAX_HISTORY = 20
last_content: dict[int, str] = load_last_content()  # 持久化，重啟後仍保留
follow_up_state: dict[int, dict] = {}   # {user_id: {"agent": "Leo", "context": "...", "active": True}}
agent_outputs: dict[int, dict] = load_agent_outputs()  # 持久化，重啟後仍保留
dreamteam_last_synthesis: dict[int, dict] = {}  # {user_id: {"question":..,"synthesis":..,"responses":..}}

REPORT_BUFFER_FILE = "/root/claude-bot/report_buffer.json"
COMPETITORS_FILE = "/root/claude-bot/competitors.json"
BRAND_INTEL_FILE = "/root/claude-bot/brand_intel.json"
pending_report_entry: dict[int, dict] = {}  # {user_id: {"agent":..,"question":..,"answer":..}}

DEFAULT_COMPETITORS = {
    "ig": [],
    "maps": [],
    "xhs": ["香港美容療程", "香港痛症"],
    "threads": [],
    "fb": [],
}

def load_competitors() -> dict:
    try:
        if os.path.exists(COMPETITORS_FILE):
            with open(COMPETITORS_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                for k in DEFAULT_COMPETITORS:
                    data.setdefault(k, [])
                return data
    except Exception:
        pass
    return dict(DEFAULT_COMPETITORS)

def save_competitors(data: dict):
    try:
        with open(COMPETITORS_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Save competitors error: {e}")

competitors_db: dict = load_competitors()

# ── 客戶智識庫 ──────────────────────────────────────────────────────────────────
DEFAULT_BRAND_INTEL: dict[str, list] = {
    "quote":     [],  # 客戶原話 / 痛點描述
    "hook":      [],  # 已驗證有效嘅開場白
    "objection": [],  # 常見反對意見 + 應對話術
    "win":       [],  # 成功案例
    "gap":       [],  # 競品空白位
}

def load_brand_intel() -> dict:
    try:
        if os.path.exists(BRAND_INTEL_FILE):
            with open(BRAND_INTEL_FILE, "r", encoding="utf-8") as f:
                data = json.load(f)
                for k in DEFAULT_BRAND_INTEL:
                    data.setdefault(k, [])
                return data
    except Exception:
        pass
    return {k: list(v) for k, v in DEFAULT_BRAND_INTEL.items()}

def save_brand_intel(data: dict):
    try:
        with open(BRAND_INTEL_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Save brand_intel error: {e}")

brand_intel_db: dict = load_brand_intel()

_INTEL_CATEGORY_ALIASES = {
    "quote": "quote", "客戶原話": "quote", "原話": "quote", "痛點": "quote",
    "hook": "hook", "開場白": "hook", "鉤子": "hook",
    "objection": "objection", "反對": "objection", "反對意見": "objection", "objections": "objection",
    "win": "win", "案例": "win", "成功案例": "win", "wins": "win",
    "gap": "gap", "空白位": "gap", "競品空白": "gap", "機會": "gap",
}

_INTEL_CATEGORY_LABELS = {
    "quote":     "📣 客戶原話",
    "hook":      "🎣 有效 Hook",
    "objection": "🛡 反對意見應對",
    "win":       "🏆 成功案例",
    "gap":       "🔍 競品空白位",
}

def _build_brand_intel_block() -> str:
    """把智識庫整合成一段 prompt 前綴，供員工使用"""
    db = brand_intel_db
    parts = []
    if db.get("quote"):
        parts.append("【客戶原話（必須整合入內容，用佢哋自己嘅字眼）】\n" +
                     "\n".join(f"• {q}" for q in db["quote"]))
    if db.get("hook"):
        parts.append("【已驗證有效 Hook（可直接使用或改編）】\n" +
                     "\n".join(f"• {h}" for h in db["hook"]))
    if db.get("objection"):
        parts.append("【常見反對意見 + 應對（預先化解）】\n" +
                     "\n".join(f"• {o}" for o in db["objection"]))
    if db.get("win"):
        parts.append("【成功案例（可引用作社交證明）】\n" +
                     "\n".join(f"• {w}" for w in db["win"]))
    if db.get("gap"):
        parts.append("【競品空白位（Stanley 可搶佔嘅角度）】\n" +
                     "\n".join(f"• {g}" for g in db["gap"]))
    if not parts:
        return ""
    return (
        "【品牌智識庫 — 以下係真實數據，製作成品時必須整合，唔可以忽略】\n\n" +
        "\n\n".join(parts)
    )

def load_report_buffer() -> dict:
    try:
        if os.path.exists(REPORT_BUFFER_FILE):
            with open(REPORT_BUFFER_FILE, "r", encoding="utf-8") as f:
                return json.load(f)
    except Exception:
        pass
    return {}

def save_report_buffer(data: dict):
    try:
        with open(REPORT_BUFFER_FILE, "w", encoding="utf-8") as f:
            json.dump(data, f, ensure_ascii=False, indent=2)
    except Exception as e:
        logger.error(f"Save report_buffer error: {e}")

report_buffer: dict = load_report_buffer()  # {str(user_id): [{"agent","question","answer","ts"}]}


# ── Web search (DuckDuckGo — free) ─────────────────────────────────────────────

def news_search(query: str, max_results: int = 5) -> str:
    """Try Apify Google first, fall back to DuckDuckGo."""
    if APIFY_AVAILABLE and APIFY_TOKEN:
        result = scrape_google(query, max_results)
        if result and not result.startswith("["):
            return result
    return web_search(query, max_results)


def web_search(query: str, max_results: int = 5) -> str:
    if not WEB_SEARCH_AVAILABLE:
        return "[Web search 不可用]"
    import time as _time
    for attempt in range(3):
        try:
            if attempt > 0:
                _time.sleep(3 * attempt)
            with DDGS() as ddgs:
                results = list(ddgs.text(query, max_results=max_results, backend="lite"))
            if results:
                out = ""
                for r in results:
                    out += f"【{r.get('title', '')}】\n{r.get('body', '')}\n來源：{r.get('href', '')}\n\n"
                return out.strip()
        except Exception as e:
            logger.warning(f"Web search attempt {attempt+1} error: {e}")
        _time.sleep(2)
    return "無搜尋結果"


# ── Google News RSS（免費、實時、唔需要任何 token）──────────────────────────────

def fetch_rss_news(feed_url: str, max_items: int = 5) -> str:
    try:
        import xml.etree.ElementTree as ET
        req = urllib.request.Request(feed_url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(req, timeout=12) as resp:
            data = resp.read().decode("utf-8", errors="replace")
        root = ET.fromstring(data)
        items = root.findall(".//item")[:max_items]
        if not items:
            return ""
        out = ""
        for item in items:
            title = (item.findtext("title") or "").strip()
            desc = re.sub(r"<[^>]+>", "", item.findtext("description") or "").strip()[:200]
            link = (item.findtext("link") or "").strip()
            out += f"【{title}】\n{desc}\n來源：{link}\n\n"
        return out.strip()
    except Exception as e:
        logger.warning(f"RSS fetch error ({feed_url[:60]}): {e}")
        return ""


def rss_hk_news() -> str:
    feeds = [
        "https://news.google.com/rss/search?q=%E9%A6%99%E6%B8%AF+%E6%96%B0%E8%81%9E&hl=zh-HK&gl=HK&ceid=HK:zh-Hant",
        "https://news.google.com/rss/search?q=%E9%A6%99%E6%B8%AF+%E8%B2%A1%E7%B6%93&hl=zh-HK&gl=HK&ceid=HK:zh-Hant",
    ]
    parts = [r for f in feeds if (r := fetch_rss_news(f, 5))]
    return "\n\n".join(parts)


def rss_ai_news() -> str:
    feeds = [
        "https://news.google.com/rss/search?q=Claude+ChatGPT+Gemini+AI+update&hl=en-US&gl=US&ceid=US:en",
    ]
    parts = [r for f in feeds if (r := fetch_rss_news(f, 5))]
    return "\n\n".join(parts)


def rss_industry_news() -> str:
    feeds = [
        "https://news.google.com/rss/search?q=%E9%A6%99%E6%B8%AF+%E7%BE%8E%E5%AE%B9+%E7%97%9B%E7%97%87&hl=zh-HK&gl=HK&ceid=HK:zh-Hant",
    ]
    parts = [r for f in feeds if (r := fetch_rss_news(f, 4))]
    return "\n\n".join(parts)


# ── Apify scraping ─────────────────────────────────────────────────────────────

def scrape_instagram(username: str, max_posts: int = 12) -> str:
    username = username.lstrip("@").strip()

    # ── 主力：Apify Instagram Scraper（最可靠）──────────────────────────────
    if APIFY_AVAILABLE and APIFY_TOKEN:
        try:
            client = ApifyClient(APIFY_TOKEN)
            run_input = {
                "directUrls": [f"https://www.instagram.com/{username}/"],
                "resultsType": "posts",
                "resultsLimit": max_posts,
                "addParentData": False,
            }
            run = client.actor("apify/instagram-scraper").call(run_input=run_input, timeout_secs=120)
            dataset_id = getattr(run, "default_dataset_id", None) or run.get("defaultDatasetId")
            items = list(client.dataset(dataset_id).iterate_items()) if dataset_id else []
            if items:
                # Try to get profile info from first item
                first = items[0]
                owner = first.get("ownerUsername") or username
                followers = first.get("followersCount") or "?"
                output = f"📊 @{owner}（{followers} followers）最新 {len(items)} 帖：\n來源：Apify IG Scraper\n\n"
                for i, item in enumerate(items[:max_posts], 1):
                    likes = item.get("likesCount") or 0
                    comments = item.get("commentsCount") or 0
                    caption = (item.get("caption") or "（無文字）")[:200]
                    ts = (item.get("timestamp") or "")[:10]
                    output += f"{i}. [{ts}] 👍{likes} 💬{comments}\n{caption}\n\n"
                return output.strip()
        except Exception as e:
            logger.warning(f"Apify IG scrape error: {e}; trying Instaloader")

    # ── Fallback 1：Instaloader ─────────────────────────────────────────────
    if INSTALOADER_AVAILABLE:
        try:
            L = instaloader.Instaloader(
                quiet=True, download_pictures=False, download_videos=False,
                download_video_thumbnails=False, download_geotags=False,
                download_comments=False, save_metadata=False,
            )
            profile = instaloader.Profile.from_username(L.context, username)
            output = f"📊 @{username}（{profile.followers:,} followers）最新 {max_posts} 帖：\n來源：Instaloader\n\n"
            for i, post in enumerate(profile.get_posts()):
                if i >= max_posts:
                    break
                caption = (post.caption or "（無文字）")[:150]
                ts = post.date_local.strftime("%Y-%m-%d")
                output += f"{i+1}. [{ts}] 👍{post.likes} 💬{post.comments}\n{caption}\n\n"
            return output.strip()
        except Exception as e:
            logger.warning(f"Instaloader error: {e}; using web search fallback")

    return (
        f"⚠️IG_FETCH_FAILED：無法直接抓取 @{username} IG。\n"
        f"Apify 及 Instaloader 均失敗。\n"
        f"可能原因：Apify credits 不足 / IG 封鎖。\n"
        f"請去 apify.com → Billing 確認 credits，或稍後再試。"
    )

def scrape_facebook(page_name: str, max_posts: int = 10) -> str:
    if not FACEBOOK_SCRAPER_AVAILABLE:
        return web_search(f"Facebook {page_name} 最新帖文 內容", 6)
    try:
        output = f"📘 Facebook @{page_name} 最新 {max_posts} 帖：\n\n"
        count = 0
        for post in fb_get_posts(page_name, pages=3, timeout=30):
            if count >= max_posts:
                break
            text = (post.get("text") or post.get("post_text") or "（無文字）")[:150]
            likes = post.get("likes") or 0
            comments = post.get("comments") or 0
            shares = post.get("shares") or 0
            t = post.get("time")
            ts = t.strftime("%Y-%m-%d") if t else ""
            output += f"{count+1}. [{ts}] 👍{likes} 💬{comments} 🔄{shares}\n{text}\n\n"
            count += 1
        if count == 0:
            return web_search(f"Facebook {page_name} 最新帖文", 6)
        return output.strip()
    except Exception as e:
        logger.error(f"Facebook scrape error: {e}")
        return web_search(f"Facebook {page_name} 最新帖文 內容", 6)


def scrape_google_maps_reviews(query: str, max_reviews: int = 20) -> str:
    """抓取 Google Maps 某商戶嘅評論（重點係差評——競品弱點）"""
    if not APIFY_AVAILABLE or not APIFY_TOKEN:
        return web_search(f"{query} Google 評論 差評 投訴", 6)
    try:
        client = ApifyClient(APIFY_TOKEN)
        run_input = {
            "searchStringsArray": [query],
            "language": "zh-HK",
            "maxReviews": max_reviews,
            "reviewsSort": "newest",
        }
        run = client.actor("compass/google-maps-reviews-scraper").call(run_input=run_input, timeout_secs=120)
        dataset_id = getattr(run, "default_dataset_id", None) or run.get("defaultDatasetId")
        items = list(client.dataset(dataset_id).iterate_items()) if dataset_id else []
        if not items:
            return web_search(f"{query} Google 評論 差評", 6)
        output = f"📍 Google Maps 評論：{query}（共 {len(items)} 條）\n\n"
        for item in items[:max_reviews]:
            stars = item.get("stars") or item.get("rating") or "?"
            text = (item.get("text") or item.get("reviewText") or "（無文字）")[:200]
            output += f"{'⭐' * int(stars) if str(stars).isdigit() else f'{stars}星'} {text}\n\n"
        return output.strip()
    except Exception as e:
        logger.warning(f"Google Maps reviews error: {e}")
        return web_search(f"{query} Google 評論 差評 投訴", 6)


def scrape_webpage(url: str) -> str:
    if not APIFY_AVAILABLE or not APIFY_TOKEN:
        return "[Apify 未設定]"
    try:
        client = ApifyClient(APIFY_TOKEN)
        run_input = {"startUrls": [{"url": url}], "maxCrawlPages": 1}
        run = client.actor("apify/web-scraper").call(run_input=run_input)
        dataset_id = getattr(run, "default_dataset_id", None) or run["defaultDatasetId"]
        items = list(client.dataset(dataset_id).iterate_items())
        if not items:
            return f"無法抓取 {url}"
        text = items[0].get("text") or items[0].get("body") or ""
        return f"【{url}】\n{text[:4000]}"
    except Exception as e:
        return f"網頁抓取失敗：{e}"


def scrape_xhs(query: str, max_posts: int = 10) -> str:
    if not APIFY_AVAILABLE or not APIFY_TOKEN:
        return "[Apify 未設定]"
    try:
        client = ApifyClient(APIFY_TOKEN)
        run_input = {"keyword": query, "maxItems": max_posts}
        run = client.actor("microworlds/xiaohongshu-scraper").call(run_input=run_input)
        dataset_id = getattr(run, "default_dataset_id", None) or run["defaultDatasetId"]
        items = list(client.dataset(dataset_id).iterate_items())
        if not items:
            return f"[小紅書：未找到「{query}」相關帖文]"
        lines = [f"📕 小紅書「{query}」搜尋結果（{len(items)} 條）：\n"]
        for i, item in enumerate(items[:max_posts], 1):
            title = item.get("title") or item.get("name") or ""
            content = item.get("description") or item.get("content") or item.get("text") or ""
            likes = item.get("likesCount") or item.get("likes") or 0
            author = item.get("authorName") or item.get("username") or ""
            url = item.get("url") or ""
            lines.append(
                f"【{i}】{title}\n"
                f"作者：{author}  ❤️ {likes}\n"
                f"{content[:200]}\n"
                f"🔗 {url}\n"
            )
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"XHS scrape error: {e}")
        return f"[小紅書抓取失敗：{e}]"


def scrape_google(query: str, max_results: int = 8) -> str:
    if not APIFY_AVAILABLE or not APIFY_TOKEN:
        return "[Apify 未設定]"
    try:
        client = ApifyClient(APIFY_TOKEN)
        run_input = {
            "queries": [query],
            "maxPagesPerQuery": 1,
            "resultsPerPage": max_results,
            "countryCode": "hk",
            "languageCode": "zh-TW",
        }
        run = client.actor("apify/google-search-scraper").call(run_input=run_input)
        dataset_id = getattr(run, "default_dataset_id", None) or run["defaultDatasetId"]
        items = list(client.dataset(dataset_id).iterate_items())
        if not items:
            return f"[Google：未找到「{query}」相關結果]"
        results = items[0].get("organicResults", []) if items else []
        if not results:
            results = items
        lines = [f"🔍 Google「{query}」搜尋結果（{len(results)} 條）：\n"]
        for i, r in enumerate(results[:max_results], 1):
            title = r.get("title") or ""
            snippet = r.get("description") or r.get("snippet") or ""
            link = r.get("url") or r.get("link") or ""
            lines.append(f"【{i}】{title}\n{snippet[:200]}\n🔗 {link}\n")
        return "\n".join(lines)
    except Exception as e:
        logger.error(f"Google scrape error: {e}")
        return f"[Google 搜尋失敗：{e}]"


def scrape_threads(username: str, max_posts: int = 12) -> str:
    if not APIFY_AVAILABLE or not APIFY_TOKEN:
        return "[Apify 未設定]"
    try:
        client = ApifyClient(APIFY_TOKEN)
        run_input = {
            "username": [username],
            "resultsLimit": max_posts,
        }
        run = client.actor("apify/threads-scraper").call(run_input=run_input)
        dataset_id = getattr(run, "default_dataset_id", None) or run["defaultDatasetId"]
        items = list(client.dataset(dataset_id).iterate_items())
        if not items:
            return f"[@{username}] Threads 抓取不到資料，可能係私人帳號或帳號名稱錯誤。"
        output = f"🧵 Threads @{username} 最新 {len(items)} 帖：\n\n"
        for i, item in enumerate(items, 1):
            likes = item.get("likeCount") or item.get("likesCount") or 0
            replies = item.get("replyCount") or item.get("repliesCount") or 0
            text = (item.get("text") or item.get("caption") or "（無文字）")[:150]
            ts = (item.get("timestamp") or item.get("createdAt") or "")[:10]
            output += f"{i}. [{ts}] ❤️{likes} 💬{replies}\n{text}\n\n"
        return output.strip()
    except Exception as e:
        logger.error(f"Threads scrape error: {e}")
        return f"Threads 抓取失敗：{e}"


# ── Voice transcription (OpenAI Whisper) ──────────────────────────────────────

def transcribe_audio(audio_bytes: bytes, filename: str = "voice.ogg") -> str:
    # faster-whisper（本地，免費，優先）
    if FASTER_WHISPER_AVAILABLE:
        try:
            import io as _io, tempfile, os as _os
            with tempfile.NamedTemporaryFile(suffix=".ogg", delete=False) as tmp:
                tmp.write(audio_bytes)
                tmp_path = tmp.name
            model = get_faster_whisper_model()
            segments, info = model.transcribe(tmp_path, language="zh", beam_size=5)
            text = " ".join(seg.text.strip() for seg in segments).strip()
            _os.unlink(tmp_path)
            return text if text else "[錄音內容為空]"
        except Exception as e:
            logger.error(f"faster-whisper error: {e}")
            # fallback to OpenAI if available
    # OpenAI Whisper fallback
    if OPENAI_AVAILABLE and OPENAI_API_KEY:
        try:
            import io as _io
            client = openai_lib.OpenAI(api_key=OPENAI_API_KEY)
            audio_file = _io.BytesIO(audio_bytes)
            audio_file.name = filename
            transcript = client.audio.transcriptions.create(model="whisper-1", file=audio_file)
            return transcript.text
        except Exception as e:
            logger.error(f"OpenAI Whisper error: {e}")
            return f"[錄音轉錄失敗：{e}]"
    return "[語音功能未啟用：請在 VPS 安裝 faster-whisper]"


# ── PowerPoint generation (python-pptx) ───────────────────────────────────────

def generate_pptx_bytes(title: str, subtitle: str, slides: list) -> bytes:
    import io as _io
    if not PPTX_AVAILABLE:
        raise ImportError("python-pptx 未安裝，請執行：pip install python-pptx")
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    DARK = RGBColor(0x1a, 0x1a, 0x2e)
    MID  = RGBColor(0x16, 0x21, 0x3e)

    # Title slide
    sl = prs.slides.add_slide(prs.slide_layouts[0])
    sl.shapes.title.text = title
    tf = sl.shapes.title.text_frame
    tf.paragraphs[0].font.name = "Microsoft YaHei"
    tf.paragraphs[0].font.size = Pt(40)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = DARK
    try:
        ph = sl.placeholders[1]
        ph.text = subtitle
        ph.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
        ph.text_frame.paragraphs[0].font.size = Pt(24)
        ph.text_frame.paragraphs[0].font.color.rgb = MID
    except Exception:
        pass

    # Content slides
    for s in slides:
        sl = prs.slides.add_slide(prs.slide_layouts[1])
        sl.shapes.title.text = s["title"]
        sl.shapes.title.text_frame.paragraphs[0].font.name = "Microsoft YaHei"
        sl.shapes.title.text_frame.paragraphs[0].font.size = Pt(32)
        sl.shapes.title.text_frame.paragraphs[0].font.bold = True
        sl.shapes.title.text_frame.paragraphs[0].font.color.rgb = MID
        try:
            body_tf = sl.placeholders[1].text_frame
            body_tf.word_wrap = True
            for i, point in enumerate(s["points"]):
                p = body_tf.paragraphs[0] if i == 0 else body_tf.add_paragraph()
                p.text = point
                p.font.name = "Microsoft YaHei"
                p.font.size = Pt(20)
                p.level = 0
        except Exception:
            pass

    buf = _io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# ── Gamma API presentation generation ─────────────────────────────────────────
# 設定步驟：gamma.app → Settings → API → 生成 API Key
# VPS 執行：export GAMMA_API_KEY='你的key' && echo 'export GAMMA_API_KEY=...' >> /root/.bashrc

def generate_gamma_pptx(title: str, content: str) -> tuple[bytes, str]:
    """
    Generate a presentation via Gamma API.
    Returns (pptx_bytes, gamma_url).
    Requires GAMMA_API_KEY env variable.
    """
    import json as _json
    if not GAMMA_API_KEY:
        raise RuntimeError("GAMMA_API_KEY 未設定")

    prompt = f"# {title}\n\n{content}"
    payload = _json.dumps({
        "inputText": prompt,
        "format": "presentation",
        "exportAs": "pptx",
        "textOptions": {"language": "zh-tw", "tone": "professional"},
        "imageOptions": {"source": "aiGenerated", "stylePreset": "photorealistic"},
    }).encode()

    req = urllib.request.Request(
        "https://api.gamma.app/v1/generate",
        data=payload,
        headers={
            "Authorization": f"Bearer {GAMMA_API_KEY}",
            "Content-Type": "application/json",
        }
    )
    with urllib.request.urlopen(req, timeout=180) as resp:
        result = _json.loads(resp.read())

    gamma_url = result.get("gammaUrl") or result.get("viewUrl") or result.get("url") or ""
    export_url = result.get("exportUrl") or ""

    pptx_bytes = b""
    if export_url:
        dl = urllib.request.Request(export_url, headers={"User-Agent": "Mozilla/5.0"})
        with urllib.request.urlopen(dl, timeout=60) as r:
            pptx_bytes = r.read()

    return pptx_bytes, gamma_url


# ── PDF generation (reportlab) ─────────────────────────────────────────────────

def generate_pdf_bytes(title: str, content: str) -> bytes:
    import io as _io
    from datetime import datetime as _dt
    if not REPORTLAB_AVAILABLE:
        raise ImportError("reportlab 未安裝，請執行：pip install reportlab")
    pdfmetrics.registerFont(UnicodeCIDFont('STSong-Light'))
    buffer = _io.BytesIO()
    doc = SimpleDocTemplate(
        buffer, pagesize=A4,
        rightMargin=2.5*cm, leftMargin=2.5*cm,
        topMargin=2.5*cm, bottomMargin=2.5*cm
    )
    title_style = ParagraphStyle(
        'CJKTitle', fontName='STSong-Light', fontSize=18,
        spaceAfter=10, alignment=TA_CENTER,
        textColor=colors.HexColor('#1a1a2e')
    )
    heading_style = ParagraphStyle(
        'CJKHeading', fontName='STSong-Light', fontSize=13,
        spaceBefore=14, spaceAfter=6,
        textColor=colors.HexColor('#16213e'), leading=20
    )
    body_style = ParagraphStyle(
        'CJKBody', fontName='STSong-Light', fontSize=11,
        spaceAfter=5, leading=19
    )
    footer_style = ParagraphStyle(
        'CJKFooter', fontName='STSong-Light', fontSize=9,
        textColor=colors.grey, alignment=TA_CENTER
    )
    story = []
    story.append(Paragraph(title, title_style))
    story.append(HRFlowable(width="100%", thickness=1, color=colors.HexColor('#cccccc')))
    story.append(Spacer(1, 0.4*cm))
    for line in content.split('\n'):
        stripped = line.strip()
        if not stripped:
            story.append(Spacer(1, 0.15*cm))
            continue
        safe = stripped.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        if stripped.startswith('#') or (len(stripped) < 35 and stripped.endswith('：')):
            clean = safe.lstrip('#').strip()
            story.append(Paragraph(clean, heading_style))
        else:
            story.append(Paragraph(safe, body_style))
    story.append(Spacer(1, 0.5*cm))
    story.append(HRFlowable(width="100%", thickness=0.5, color=colors.HexColor('#dddddd')))
    date_str = _dt.now().strftime("%Y年%m月%d日")
    story.append(Paragraph(f"生成日期：{date_str}", footer_style))
    doc.build(story)
    return buffer.getvalue()


# ── Image generation (Pollinations.ai — free) ──────────────────────────────────

def generate_image_sync(prompt: str) -> bytes:
    encoded = urllib.parse.quote(prompt)
    url = f"https://image.pollinations.ai/prompt/{encoded}?width=1080&height=1080&nologo=true&model=flux"
    req = urllib.request.Request(url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(req, timeout=90) as resp:
        return resp.read()


def generate_dalle_image(prompt: str, size: str = "1024x1024") -> bytes:
    """Generate image via DALL-E 3 (OpenAI). Falls back to Pollinations if unavailable."""
    if not OPENAI_AVAILABLE or not OPENAI_API_KEY:
        return generate_image_sync(prompt)
    import io as _io
    client = openai_lib.OpenAI(api_key=OPENAI_API_KEY)
    response = client.images.generate(
        model="dall-e-3",
        prompt=prompt,
        size=size,
        quality="standard",
        n=1,
    )
    image_url = response.data[0].url
    dl_req = urllib.request.Request(image_url, headers={"User-Agent": "Mozilla/5.0"})
    with urllib.request.urlopen(dl_req, timeout=60) as resp:
        return resp.read()


def generate_dalle_prompts(user_desc: str) -> str:
    user_msg = (
        f"用戶需求：{user_desc}\n\n"
        "業務背景：香港美容（護膚、美容療程）及痛症管理銷售，IG 行銷為主。\n\n"
        "請生成 3 個不同風格的英文 DALL-E 3 prompt，每個方向不同：\n"
        "1. 🏥 專業醫療風（乾淨、信任感、白色基調）\n"
        "2. 💆 生活方式風（溫暖、真實感、自然光）\n"
        "3. ✨ 高端美學風（奢華、精緻、品牌感）\n\n"
        "每個 prompt 要包含：主體描述、風格、光線、構圖、色調、--ar 9:16（IG 規格）。\n"
        "直接輸出 3 個 prompt，唔需要額外解釋。"
    )
    return run_with_system(
        "你係一位頂尖 AI 圖像設計師，專門幫人優化 DALL-E 3 / ChatGPT Image 的 prompt。",
        user_msg,
    )


# ── Core Claude runners ────────────────────────────────────────────────────────

def run_claude(prompt: str, timeout: int = 180, model: str = None) -> str:
    env = {**os.environ, "HOME": "/root"}
    env.pop("ANTHROPIC_API_KEY", None)  #用 OAuth 登入，唔用 API key
    cmd = [CLAUDE_PATH, "-p", prompt, "--allowedTools", ""]
    if model:
        cmd += ["--model", model]
    result = subprocess.run(
        cmd,
        capture_output=True, text=True, timeout=timeout,
        env=env,
    )
    track_call()
    combined_out = (result.stdout or "") + (result.stderr or "")
    if any(kw in combined_out.lower() for kw in ["rate limit", "usage limit", "quota", "too many requests", "overloaded", "credit balance"]):
        return "[QUOTA_EXCEEDED]"
    if result.returncode == 0 and result.stdout.strip():
        return result.stdout.strip()
    logger.error(f"Claude error: stderr={result.stderr[:200]} stdout={result.stdout[:200]}")
    return f"[出錯：{result.stderr[:100]}]"


def run_claude_sdk(messages: list, model: str = MODEL_FAST, max_tokens: int = 4096, system: str = None) -> str:
    try:
        client = anthropic.Anthropic()
        kwargs = dict(model=model, max_tokens=max_tokens, messages=messages)
        if system:
            kwargs['system'] = system
        response = client.messages.create(**kwargs)
        track_call()
        return response.content[0].text
    except Exception as e:
        err_str = str(e).lower()
        if any(kw in err_str for kw in ["rate limit", "usage limit", "quota", "too many requests", "overloaded", "credit balance"]):
            return "[QUOTA_EXCEEDED]"
        logger.error(f"Claude SDK error: {e}")
        return f"[API 錯誤：{e}]"


def run_with_system(system_prompt: str, user_message: str, model: str = MODEL_FAST, timeout: int = None) -> str:
    combined = f"{system_prompt}\n\n---\n{user_message}"
    if timeout is None:
        timeout = 300 if model == MODEL_HEAVY else 180
    return run_claude(combined, timeout=timeout, model=model)


def extract_file_content(file_bytes: bytes, mime: str, filename: str = "") -> str:
    import tempfile as _tf, os as _os

    if mime == "application/pdf":
        try:
            import pdfplumber, io as _io
            with pdfplumber.open(_io.BytesIO(file_bytes)) as pdf:
                text = "\n".join(p.extract_text() or "" for p in pdf.pages)
            if text.strip():
                return text[:10000]
        except Exception:
            pass
        # fallback: base64 via CLI
        with _tf.NamedTemporaryFile(suffix=".pdf", delete=False) as f:
            f.write(file_bytes)
            tmp = f.name
        try:
            result = run_claude(f"請完整提取並整理呢份PDF文件嘅所有內容，包括標題、段落、數據、重點。保持原有結構。\n\n[附件：{tmp}]")
        finally:
            _os.unlink(tmp)
        return result

    elif mime.startswith("image/"):
        suffix = ".jpg" if "jpeg" in mime else f".{mime.split('/')[-1]}"
        with _tf.NamedTemporaryFile(suffix=suffix, delete=False) as f:
            f.write(file_bytes)
            tmp = f.name
        prompt = (
            "請仔細分析呢張圖片。\n"
            "如果係手寫文字/筆記：完整轉錄所有文字內容，保持原有結構。\n"
            "如果係截圖/文件：提取所有文字同數據。\n"
            "如果係產品/場景圖：詳細描述內容、風格、可用於行銷嘅重點。\n"
            "廣東話回答，盡量完整。"
        )
        try:
            result = subprocess.run(
                [CLAUDE_PATH, "-p", prompt, tmp],
                capture_output=True, text=True, timeout=120,
                env={**os.environ, "HOME": "/root"},
            )
            if result.returncode == 0 and result.stdout.strip():
                return result.stdout.strip()
        finally:
            _os.unlink(tmp)
        return "[圖片分析失敗]"

    elif mime in ("text/plain", "text/csv", "text/markdown", "application/json") or \
         any(filename.endswith(ext) for ext in (".txt", ".csv", ".md", ".json")):
        return file_bytes.decode("utf-8", errors="ignore")[:10000]

    elif mime in ("application/vnd.openxmlformats-officedocument.wordprocessingml.document",) or \
         filename.endswith(".docx"):
        try:
            import docx
            import io
            doc = docx.Document(io.BytesIO(file_bytes))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            return "[需要安裝 python-docx 才能讀取 Word 文件]"
        except Exception as e:
            return f"[讀取 Word 文件出錯：{e}]"

    return ""


# ── Agent Avatar Cards ────────────────────────────────────────────────────────
AGENT_AVATAR_DIR = "/root/claude-bot/avatars"

AGENT_AVATAR_DATA = {
    "Amy":   {"bg": (255, 200, 220), "hair": (220,  60, 120), "clothes": (160,  80, 180), "role": "首席秘書"},
    "Anna":  {"bg": (255, 235, 190), "hair": (200, 110,  20), "clothes": (230, 140,  30), "role": "內容創作師"},
    "Leo":   {"bg": (195, 220, 255), "hair": ( 40,  60, 160), "clothes": ( 60, 100, 210), "role": "情報分析師"},
    "Kai":   {"bg": (195, 245, 210), "hair": ( 20, 140,  55), "clothes": ( 35, 170,  75), "role": "AI 情報員"},
    "Toxic": {"bg": (255, 255, 190), "hair": (190, 150,   0), "clothes": (215, 175,   0), "role": "自動化工程師"},
    "Small": {"bg": (220, 200, 255), "hair": (100,  30, 165), "clothes": (125,  55, 185), "role": "商業策略官"},
    "Tony":  {"bg": (255, 210, 200), "hair": (155,  30,  20), "clothes": (195,  55,  40), "role": "轉化專員"},
    "Rex":   {"bg": (255, 225, 185), "hair": (155,  75,   0), "clothes": (195, 100,  10), "role": "廣告投放"},
    "Mia":   {"bg": (195, 245, 245), "hair": (  0, 115, 115), "clothes": (  0, 155, 155), "role": "數據分析師"},
}

AGENT_EMOJI_LABEL = {
    "Amy": "👩‍💼", "Anna": "🎨", "Leo": "📊", "Kai": "🤖",
    "Toxic": "⚡", "Small": "🧠", "Tony": "🤝", "Rex": "📢", "Mia": "📈",
}

def _make_agent_avatar(agent_name: str) -> bytes:
    """Generate pixel-art agent card: draw at 1/4 res, scale up 4× with NEAREST."""
    from PIL import Image, ImageDraw, ImageFont
    import io as _io

    SCALE = 4
    W, H = 50, 68
    d = AGENT_AVATAR_DATA.get(agent_name, AGENT_AVATAR_DATA["Amy"])
    bg, hair, clothes = d["bg"], d["hair"], d["clothes"]
    SKIN  = (252, 218, 168)
    EYE   = (40, 28, 16)
    MOUTH = (210, 75, 75)
    SHOE  = (55, 40, 25)
    PANTS = tuple(max(0, c - 40) for c in clothes)
    DARK  = tuple(max(0, c - 28) for c in bg)
    DARK_C = tuple(max(0, c - 28) for c in clothes)

    img = Image.new("RGB", (W, H), bg)
    draw = ImageDraw.Draw(img)
    draw.rectangle([0, 0, W - 1, H - 1], outline=DARK, width=1)

    cx = W // 2  # 25

    # Hair blob
    draw.ellipse([cx - 9, 2, cx + 9, 14], fill=hair)
    # Face oval
    draw.ellipse([cx - 8, 7, cx + 8, 21], fill=SKIN)
    # Hair sides
    draw.rectangle([cx - 9, 8, cx - 7, 15], fill=hair)
    draw.rectangle([cx + 7, 8, cx + 9, 15], fill=hair)
    # Eyes
    draw.point([cx - 4, 13], fill=EYE)
    draw.point([cx + 4, 13], fill=EYE)
    # Mouth
    draw.rectangle([cx - 2, 17, cx + 2, 17], fill=MOUTH)

    # Body
    by = 23
    draw.rectangle([cx - 7, by, cx + 7, by + 13], fill=clothes)
    # Arms
    draw.rectangle([cx - 10, by + 1, cx - 8, by + 10], fill=clothes)
    draw.rectangle([cx + 8,  by + 1, cx + 10, by + 10], fill=clothes)
    # Hands
    draw.ellipse([cx - 12, by + 9, cx - 7, by + 13], fill=SKIN)
    draw.ellipse([cx + 7,  by + 9, cx + 12, by + 13], fill=SKIN)
    # Collar
    draw.rectangle([cx - 3, by, cx + 3, by + 5], fill=DARK_C)

    # Legs
    ly = by + 14
    draw.rectangle([cx - 6, ly, cx - 2, ly + 11], fill=PANTS)
    draw.rectangle([cx + 2, ly, cx + 6, ly + 11], fill=PANTS)
    # Shoes
    sy = ly + 10
    draw.rectangle([cx - 8, sy, cx - 1, sy + 3], fill=SHOE)
    draw.rectangle([cx + 1, sy, cx + 8, sy + 3], fill=SHOE)

    # Scale up → pixel art effect
    img = img.resize((W * SCALE, H * SCALE), Image.NEAREST)
    draw = ImageDraw.Draw(img)
    FW, FH = W * SCALE, H * SCALE

    # Footer bar
    footer_bg = tuple(max(0, c - 50) for c in bg)
    draw.rectangle([0, FH - 46, FW, FH], fill=footer_bg)

    try:
        name_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf", 16)
        role_font = ImageFont.truetype("/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf", 11)
    except Exception:
        name_font = role_font = ImageFont.load_default()

    draw.text((FW // 2, FH - 32), agent_name,      font=name_font, fill=(255, 255, 255), anchor="mm")
    draw.text((FW // 2, FH - 14), d.get("role",""), font=role_font, fill=(200, 200, 200), anchor="mm")

    buf = _io.BytesIO()
    img.save(buf, "PNG")
    return buf.getvalue()


async def send_with_avatar(update: Update, agent_name: str, text: str):
    """Send pixel-art agent card, then the text response."""
    os.makedirs(AGENT_AVATAR_DIR, exist_ok=True)
    avatar_path = os.path.join(AGENT_AVATAR_DIR, f"{agent_name}.png")
    loop = asyncio.get_event_loop()

    if not os.path.exists(avatar_path):
        img_bytes = await loop.run_in_executor(executor, _make_agent_avatar, agent_name)
        with open(avatar_path, "wb") as f:
            f.write(img_bytes)

    emoji = AGENT_EMOJI_LABEL.get(agent_name, "")
    role  = AGENT_AVATAR_DATA.get(agent_name, {}).get("role", "")
    caption = f"{emoji} *{agent_name}* — {role}"
    try:
        with open(avatar_path, "rb") as f:
            await update.message.reply_photo(photo=f, caption=caption, parse_mode="Markdown")
    except Exception as e:
        logger.warning(f"Avatar send failed ({agent_name}): {e}")

    await send_long(update, text)


_INTEL_INJECT_AGENTS = {"Anna", "Tony", "Rex", "Leo", "Small"}

def agent_call(agent_name: str, task: str) -> tuple[str, str]:
    if agent_name in _INTEL_INJECT_AGENTS:
        intel_block = _build_brand_intel_block()
        if intel_block:
            task = intel_block + "\n\n---\n\n" + task
    user_msg = f"Stanley 交俾你嘅任務：\n{task}\n\n請根據你嘅角色同專長，提供最完整、最具體、可直接使用嘅成品，唔好省略任何重要細節。"
    model = AGENT_MODELS.get(agent_name, MODEL_FAST)
    timeout = 300 if model == MODEL_HEAVY else 180
    return agent_name, run_with_system(AGENT_PROMPTS[agent_name], user_msg, model, timeout)


# ── Commands ───────────────────────────────────────────────────────────────────

async def cmd_start(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    conversation_history[ALLOWED_USER_ID] = []
    save_history(conversation_history)
    last_content[ALLOWED_USER_ID] = ""
    save_last_content_to_disk(last_content)
    search_st = "✅" if WEB_SEARCH_AVAILABLE else "❌"
    apify_st  = "✅" if (APIFY_AVAILABLE and APIFY_TOKEN) else "❌"
    voice_st = "✅" if (OPENAI_AVAILABLE and OPENAI_API_KEY) else "❌"
    pdf_st   = "✅" if REPORTLAB_AVAILABLE else "❌"
    await update.message.reply_text(
        "Stanley，我係 Amy v4.0 ✅\n\n"
        "全體就位：Anna / Leo / Kai / Toxic / Small / Tony / Rex / Mia\n\n"
        "━━ 支援嘅輸入 ━━\n"
        "📄 PDF  🖼 圖片  ✍️ 手寫圖  📝 Word/TXT/CSV\n"
        "🎙️ 語音訊息  💬 文字指令\n\n"
        "━━ 成品輸出 ━━\n"
        "/pdf [描述] → 生成 PDF 文件\n"
        "/slides [描述] → 生成 PowerPoint\n"
        "/landingpage [描述] → 生成 Landing Page\n"
        "/caption [描述] → IG/小紅書 Caption\n"
        "/reel [主題] → Reels 腳本\n"
        "/copy [需求] → 廣告文案\n\n"
        "━━ 情報工具 ━━\n"
        "/imagine [描述] → AI 圖片 + Prompts\n"
        "/scrape ig @帳號 → IG 競品分析\n"
        "/scrape web [URL] → 抓取網頁內容\n"
        "/xhs [關鍵詞] → 小紅書帖文抓取\n"
        "/research [關鍵詞] → Google 情報搜尋\n"
        "/report → 即時市場 + AI 情報\n"
        "/dreamteam [問題] → 8位教練分析（唔帶問題→審閱最新內容）\n"
        "/dreamteampdf → Dream Team 分析生成 PDF\n"
        "/myreport → 睇所有批准嘅 Q&A\n"
        "/reportpdf → Report Q&A 生成 PDF\n\n"
        f"🔍 網絡搜尋：{search_st}  📡 Apify：{apify_st}\n"
        f"🎙️ 語音轉錄：{voice_st}  📄 PDF：{pdf_st}\n"
        "每日 8:45 自動發送情報簡報（頭條新聞 + 行業要聞 + AI 資訊）。\n\n"
        "有咩吩咐？"
    )


async def cmd_clear(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    conversation_history[ALLOWED_USER_ID] = []
    save_history(conversation_history)
    last_content[ALLOWED_USER_ID] = ""
    save_last_content_to_disk(last_content)
    await update.message.reply_text("記憶清除，重新開始。")


async def cmd_status(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    h = conversation_history.get(ALLOWED_USER_ID, [])
    search_st = "✅ 啟用（真實新聞）" if WEB_SEARCH_AVAILABLE else "❌ 需安裝 duckduckgo-search"
    apify_st  = "✅ 啟用" if (APIFY_AVAILABLE and APIFY_TOKEN) else "❌ 需設定 APIFY_API_TOKEN"
    if FASTER_WHISPER_AVAILABLE:
        voice_st = "✅ 啟用（本地 Whisper，免費）"
    elif OPENAI_AVAILABLE and OPENAI_API_KEY:
        voice_st = "✅ 啟用（OpenAI Whisper）"
    else:
        voice_st = "❌ 未啟用（pip install faster-whisper）"
    pdf_st   = "✅ 啟用" if REPORTLAB_AVAILABLE else "❌ 需安裝 reportlab"
    dalle_st = "✅ DALL-E 3（真實生成）" if (OPENAI_AVAILABLE and OPENAI_API_KEY) else "⚡ Pollinations.ai（免費）"
    gamma_st = "✅ Gamma AI（專業簡報）" if GAMMA_API_KEY else "❌ 未設定（用 python-pptx）\n   設定：export GAMMA_API_KEY='你的key'"
    pptx_st  = "✅ 啟用" if PPTX_AVAILABLE else "❌ 需安裝 python-pptx"
    await update.message.reply_text(
        f"系統狀態：正常\n"
        f"版本：v4.2\n"
        f"對話記憶：{len(h)} 條（持久化✅）\n\n"
        f"━━ AI 分工模型 ━━\n"
        f"Amy（指揮）：Haiku ⚡\n"
        f"Anna（文案）：Sonnet ✍️\n"
        f"Small（策略）：Opus 🧠\n"
        f"Kai/Mia（情報）：Haiku ⚡\n"
        f"其他員工：Sonnet\n\n"
        f"━━ 外部 AI 工具 ━━\n"
        f"🖼 圖片生成：{dalle_st}\n"
        f"📊 簡報生成：{gamma_st}\n"
        f"📄 PDF：{pdf_st}\n"
        f"🎙️ 語音轉錄：{voice_st}\n"
        f"🔍 網絡搜尋：{search_st}\n"
        f"📡 RSS 實時新聞：✅ 免費啟用（Google News）\n"
        f"📡 Apify 爬蟲：{apify_st}"
    )


async def cmd_export(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    from datetime import datetime
    import io
    h = conversation_history.get(ALLOWED_USER_ID, [])
    lines = [
        f"Stanley Amy Bot — 對話記錄匯出",
        f"匯出時間：{datetime.now().strftime('%Y-%m-%d %H:%M:%S')} (HKT)",
        f"總共 {len(h)} 條記錄（最近 50 條）",
        "=" * 50,
        "",
    ]
    for i, msg in enumerate(h, 1):
        role = "👤 你" if msg.get("role") == "user" else "🤖 Amy"
        content = msg.get("content", "")
        if isinstance(content, list):
            content = " ".join(c.get("text", "") for c in content if isinstance(c, dict))
        lines.append(f"[{i}] {role}")
        lines.append(content.strip())
        lines.append("")
    lc = last_content.get(ALLOWED_USER_ID, "")
    if lc:
        lines += ["=" * 50, "📄 上次處理嘅內容（節錄）：", lc[:500] + ("..." if len(lc) > 500 else ""), ""]
    usage = _daily_usage.get("calls", 0)
    lines.append(f"今日 API 使用次數：{usage}")
    text = "\n".join(lines)
    buf = io.BytesIO(text.encode("utf-8"))
    buf.name = f"amy_history_{datetime.now().strftime('%Y%m%d_%H%M')}.txt"
    await update.message.reply_document(document=buf, filename=buf.name, caption=f"匯出完成，共 {len(h)} 條對話記錄。")


async def cmd_imagine(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    prompt = " ".join(context.args) if context.args else ""
    if not prompt:
        dalle_on = OPENAI_AVAILABLE and bool(OPENAI_API_KEY)
        await update.message.reply_text(
            "用法：/imagine [描述]\n\n"
            "例子：\n"
            "• /imagine 香港現代痛症診所廣告，專業溫馨\n"
            "• /imagine 美容護膚產品平面廣告，白色背景，高質感\n"
            "• /imagine IG封面圖，美容主題，粉色系\n\n"
            f"圖片引擎：{'✅ DALL-E 3（真實生成）' if dalle_on else '⚡ Pollinations.ai + DALL-E 3 Prompts'}"
        )
        return

    use_dalle = OPENAI_AVAILABLE and bool(OPENAI_API_KEY)
    img_source = "DALL-E 3 (GPT-4o)" if use_dalle else "Pollinations.ai"
    await update.message.reply_text(f"🎨 Anna 處理中（{img_source} 生成圖片，約30-60秒）...")

    loop = asyncio.get_event_loop()
    if use_dalle:
        img_bytes = await loop.run_in_executor(executor, generate_dalle_image, prompt)
        if isinstance(img_bytes, Exception):
            await update.message.reply_text(f"⚠️ 圖片生成失敗：{img_bytes}")
        else:
            await update.message.reply_photo(photo=img_bytes, caption=f"🎨 {img_source}\n描述：{prompt}")
    else:
        img_task = loop.run_in_executor(executor, generate_image_sync, prompt)
        dalle_task = loop.run_in_executor(executor, generate_dalle_prompts, prompt)
        img_bytes, dalle_prompts = await asyncio.gather(img_task, dalle_task, return_exceptions=True)
        if isinstance(img_bytes, Exception):
            await update.message.reply_text(f"⚠️ 圖片生成失敗：{img_bytes}")
        else:
            await update.message.reply_photo(photo=img_bytes, caption=f"🎨 {img_source}\n描述：{prompt}")
        if not isinstance(dalle_prompts, Exception):
            await send_long(update, f"✨ DALL-E 3 優化 Prompts（貼入 ChatGPT → GPT-4o）：\n\n{dalle_prompts}")


async def cmd_scrape(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return

    args = context.args
    if not args:
        await update.message.reply_text(
            "用法：\n"
            "• /scrape ig @帳號 — 抓取 IG 競品最新帖文 + Leo 分析\n"
            "• /scrape web [URL] — 抓取網頁內容\n\n"
            "例：/scrape ig @beautysignaturehk"
        )
        return

    subcmd = args[0].lower()
    loop = asyncio.get_event_loop()

    def _scrape_failed(r: str) -> bool:
        return any(kw in r for kw in ["失敗", "未設定", "不到資料", "Error", "error", "抓取失敗", "⚠️", "FAILED"])

    def _save_analysis(analysis: str, label: str):
        """分析完成後儲存到 last_content + history，方便後續生成用"""
        last_content[ALLOWED_USER_ID] = analysis[:8000]
        save_last_content_to_disk(last_content)
        h = conversation_history.setdefault(ALLOWED_USER_ID, [])
        h.append({"role": "assistant", "content": f"[{label}]\n{analysis}"})
        save_history(conversation_history)

    _next_steps = (
        "\n💡 分析已儲存！你可以直接用：\n"
        "• /landingpage 根據以上分析生成落地頁\n"
        "• /pdf 生成競品分析報告\n"
        "• /slides 生成競品分析簡報\n"
        "• 或問我：「根據以上分析幫我設計新產品」"
    )

    if subcmd == "ig" and len(args) >= 2:
        username = args[1].lstrip("@")
        await update.message.reply_text(f"📡 抓取 IG @{username} 中（約30-60秒）...")
        result = await loop.run_in_executor(executor, scrape_instagram, username)
        await send_long(update, result)
        if _scrape_failed(result):
            await update.message.reply_text(
                "❌ IG 數據抓取失敗，Leo 分析已取消（節省 tokens）。\n"
                "請去 apify.com → Billing 確認 credits 充足後再試。"
            )
            return
        leo_task = (
            f"以下係競品 IG 帳號 @{username} 嘅最新帖文數據，請做完整競品分析：\n\n{result}\n\n"
            f"輸出：① 帳號整體定位 ② 爆款帖文規律 ③ 對 Stanley 美容/痛症業務嘅具體啟示 ④ 可以抄嘅策略"
        )
        _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
        await send_with_avatar(update, "Leo", analysis)
        _save_analysis(analysis, f"IG @{username} 競品分析")
        await update.message.reply_text(_next_steps)

    elif subcmd == "web" and len(args) >= 2:
        url = args[1]
        await update.message.reply_text(f"📡 抓取 {url} 中...")
        result = await loop.run_in_executor(executor, scrape_webpage, url)
        await send_long(update, result)
        if _scrape_failed(result):
            await update.message.reply_text("❌ 網頁抓取失敗，Leo 分析取消（節省 tokens）。")
            return
        leo_task = f"分析以下網頁內容，提出對 Stanley 業務有用嘅洞察：\n\n{result[:3000]}"
        _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
        await send_with_avatar(update, "Leo", analysis)
        _save_analysis(analysis, f"網頁分析 {url}")
        await update.message.reply_text(_next_steps)

    elif subcmd == "xhs" and len(args) >= 2:
        query = " ".join(args[1:])
        await update.message.reply_text(f"📕 抓取小紅書「{query}」中（約30-60秒）...")
        result = await loop.run_in_executor(executor, scrape_xhs, query)
        await send_long(update, result)
        if _scrape_failed(result):
            await update.message.reply_text("❌ 小紅書抓取失敗，Anna 分析取消（節省 tokens）。\n請確認 Apify credits 充足。")
            return
        anna_task = (
            f"以下係小紅書「{query}」嘅帖文內容，請分析趨勢：\n\n{result[:3000]}\n\n"
            f"輸出：① 爆款標題規律 ② 熱門話題角度 ③ Stanley 可以用嘅3個文案方向"
        )
        _, analysis = await loop.run_in_executor(executor, agent_call, "Anna", anna_task)
        await send_with_avatar(update, "Anna", analysis)
        _save_analysis(analysis, f"小紅書「{query}」趨勢分析")
        await update.message.reply_text(_next_steps)

    elif subcmd == "news" and len(args) >= 2:
        query = " ".join(args[1:])
        await update.message.reply_text(f"🔍 搜尋「{query}」新聞中...")
        result = await loop.run_in_executor(executor, scrape_google, query)
        await send_long(update, result)
        if _scrape_failed(result):
            await update.message.reply_text("❌ 新聞搜尋失敗，Leo 分析取消（節省 tokens）。")
            return
        leo_task = f"以下係「{query}」嘅最新搜尋結果，請分析對 Stanley 業務嘅影響同機會：\n\n{result[:3000]}"
        _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
        await send_with_avatar(update, "Leo", analysis)
        _save_analysis(analysis, f"新聞搜尋「{query}」分析")
        await update.message.reply_text(_next_steps)

    elif subcmd == "threads" and len(args) >= 2:
        username = args[1].lstrip("@")
        await update.message.reply_text(f"🧵 抓取 Threads @{username} 中（約30-60秒）...")
        result = await loop.run_in_executor(executor, scrape_threads, username)
        await send_long(update, result)
        if _scrape_failed(result):
            await update.message.reply_text("❌ Threads 抓取失敗，Leo 分析取消（節省 tokens）。\n請確認 Apify credits 充足。")
            return
        leo_task = (
            f"以下係競品 Threads 帳號 @{username} 嘅最新內容，請做競品分析：\n\n{result}\n\n"
            f"輸出：① 帳號內容定位 ② 爆款帖文規律 ③ 對 Stanley 業務嘅具體啟示 ④ 可以借鑒嘅策略"
        )
        _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
        await send_with_avatar(update, "Leo", analysis)
        _save_analysis(analysis, f"Threads @{username} 競品分析")
        await update.message.reply_text(_next_steps)

    elif subcmd == "fb" and len(args) >= 2:
        page = args[1].lstrip("@")
        await update.message.reply_text(f"📘 抓取 Facebook @{page} 中...")
        result = await loop.run_in_executor(executor, scrape_facebook, page)
        await send_long(update, result)
        leo_task = (
            f"以下係競品 Facebook 專頁 @{page} 嘅最新帖文，請做競品分析：\n\n{result}\n\n"
            f"輸出：① 專頁內容定位 ② 爆款帖文規律 ③ 對 Stanley 業務嘅具體啟示 ④ 可以借鑒嘅策略"
        )
        _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
        await send_with_avatar(update, "Leo", analysis)
        _save_analysis(analysis, f"Facebook @{page} 競品分析")
        await update.message.reply_text(_next_steps)

    else:
        await update.message.reply_text(
            "用法：\n"
            "• /scrape ig @帳號 — IG 競品分析\n"
            "• /scrape fb @專頁 — Facebook 競品分析\n"
            "• /scrape threads @帳號 — Threads 競品分析\n"
            "• /scrape web [URL] — 網頁內容\n"
            "• /scrape xhs [關鍵詞] — 小紅書帖文\n"
            "• /scrape news [關鍵詞] — Google 新聞搜尋"
        )


async def cmd_xhs(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    query = " ".join(context.args) if context.args else ""
    if not query:
        await update.message.reply_text(
            "用法：/xhs [關鍵詞]\n\n"
            "例子：\n"
            "• /xhs 痛症舒緩\n"
            "• /xhs 美白護膚 香港"
        )
        return
    if not (APIFY_AVAILABLE and APIFY_TOKEN):
        await update.message.reply_text("⚠️ Apify 未設定，請先設定 APIFY_API_TOKEN")
        return
    await update.message.reply_text(f"📕 抓取小紅書「{query}」中（約30-60秒）...")
    loop = asyncio.get_event_loop()
    result = await loop.run_in_executor(executor, scrape_xhs, query)
    await send_long(update, result)
    anna_task = (
        f"以下係小紅書「{query}」嘅帖文，請分析趨勢同文案方向：\n\n{result[:3000]}\n\n"
        f"輸出：① 爆款標題規律 ② 熱門話題角度 ③ Stanley 可以直接用嘅3個文案方向（附範例）"
    )
    _, analysis = await loop.run_in_executor(executor, agent_call, "Anna", anna_task)
    await send_long(update, f"🎨 Anna 分析：\n{analysis}")


async def cmd_research(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    query = " ".join(context.args) if context.args else ""
    if not query:
        await update.message.reply_text(
            "用法：/research [關鍵詞]\n\n"
            "例子：\n"
            "• /research 香港痛症治療市場\n"
            "• /research 競品美容中心優惠 2025"
        )
        return
    if not (APIFY_AVAILABLE and APIFY_TOKEN):
        await update.message.reply_text("⚠️ Apify 未設定，請先設定 APIFY_API_TOKEN")
        return
    await update.message.reply_text(f"🔍 搜尋「{query}」中...")
    loop = asyncio.get_event_loop()
    result = await loop.run_in_executor(executor, scrape_google, query)
    await send_long(update, result)
    leo_task = (
        f"以下係「{query}」嘅 Google 搜尋結果，請分析對 Stanley 業務嘅影響同機會：\n\n{result[:3000]}\n\n"
        f"輸出：① 核心洞察（3點）② 競品動態 ③ 可立即執行嘅行動建議"
    )
    _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
    await send_long(update, f"📊 Leo 分析：\n{analysis}")


async def cmd_landingpage(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return

    description = " ".join(context.args) if context.args else ""
    if not description:
        cached = last_content.get(ALLOWED_USER_ID, "")
        history = conversation_history.get(ALLOWED_USER_ID, [])
        recent = "\n".join(
            f"{'Stanley' if m['role']=='user' else 'Amy'}：{m['content'][:800]}"
            for m in history[-4:]
        )
        if recent or cached:
            context_text = ""
            if recent:
                context_text += f"最近分析：\n{recent}\n\n"
            if cached and not recent:
                context_text += f"參考內容：\n{cached[:3000]}\n\n"
            description = f"根據以下競品分析，生成針對性 Landing Page：\n\n{context_text}"
        else:
            await update.message.reply_text(
                "用法：/landingpage [描述]\n\n"
                "例子：\n"
                "• /landingpage 痛症治療，目標30-60歲，賣點免費初診\n"
                "• /landingpage 美容護膚課程，25-40歲女性，天然成分\n\n"
                "或先用 /scrape 分析競品，再直接 /landingpage 即可自動生成。"
            )
            return

    await update.message.reply_text("🎨 Anna 設計緊，請稍等（約30秒）...")

    user_msg = (
        f"Stanley 需要一個 Landing Page。\n描述：{description}\n\n"
        f"生成完整 HTML Landing Page（單一文件）：\n"
        f"- 用 Tailwind CSS CDN\n"
        f"- 包含：Hero、服務介紹、痛點解決、社會認同（testimonial）、CTA\n"
        f"- 繁體中文，廣東話語氣，手機優先\n"
        f"- 只輸出純 HTML，唔需要解釋"
    )

    loop = asyncio.get_event_loop()
    html_content = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Anna'], user_msg)

    html_match = re.search(r'```(?:html)?\s*(<!DOCTYPE.*?</html>)\s*```', html_content, re.DOTALL | re.IGNORECASE)
    if html_match:
        html_content = html_match.group(1)

    with tempfile.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as f:
        f.write(html_content)
        tmp_path = f.name

    try:
        with open(tmp_path, 'rb') as f:
            await update.message.reply_document(document=f, filename='landing_page.html',
                                                 caption="🎨 Anna 完成！用瀏覽器打開預覽。")
    finally:
        os.unlink(tmp_path)


async def cmd_report(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    await update.message.reply_text("📊 生成緊情報簡報，請稍等...")
    await send_daily_report(context.application)


async def cmd_dreamteam(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return

    question = " ".join(context.args) if context.args else ""
    pdf_mode = not question

    if pdf_mode:
        content_ctx = last_content.get(ALLOWED_USER_ID, "")
        if not content_ctx:
            await update.message.reply_text(
                "⚠️ 未找到最新分析或 PDF 內容。\n\n"
                "請先讓員工做分析，或生成 PDF，Dream Team 先可以幫你審閱。\n"
                "或者：/dreamteam [你嘅問題] 直接問教練。"
            )
            return
        display_question = "審閱你最新嘅分析內容"
        coach_user_msg = (
            f"以下係 Stanley 最新嘅分析/文件內容，請根據你嘅專長，"
            f"就呢份內容提供具體改進意見、盲點、或行動建議：\n\n{content_ctx[:4000]}"
        )
        await update.message.reply_text(
            "🏆 Dream Team 審閱模式\n針對你最新嘅分析內容，8位教練並行俾意見（約60-90秒）..."
        )
    else:
        display_question = question
        coach_user_msg = f"Stanley 的問題：{question}\n\n請根據你的專長，提供具體、可執行的建議。"
        await update.message.reply_text(
            f"🏆 Dream Team 集結中...\n問題：{question}\n\n8位教練並行分析（約60-90秒）..."
        )

    loop = asyncio.get_event_loop()

    def coach_call(name: str, system_prompt: str) -> tuple[str, str]:
        model = DREAM_TEAM_MODELS.get(name, MODEL_FAST)
        return name, run_with_system(system_prompt, coach_user_msg, model)

    tasks = [loop.run_in_executor(executor, coach_call, n, p) for n, p in DREAM_TEAM_PROMPTS.items()]
    results = list(await asyncio.gather(*tasks))

    all_responses = ""
    for name, reply in results:
        emoji = DREAM_TEAM_EMOJI.get(name, "🏆")
        msg = f"{emoji} {name}：\n{reply}"
        all_responses += msg + "\n\n"
        await send_long(update, msg)

    synthesis_user = (
        f"以下係8位世界頂尖商業教練嘅意見：\n\n"
        f"議題：{display_question}\n\n{all_responses}\n\n"
        f"請用廣東話整合核心建議，找出共識、點出分歧、提出最優先嘅3個行動點。格式清晰。"
    )
    synthesis = await loop.run_in_executor(
        executor, run_with_system,
        "你係一個商業策略整合師，廣東話輸出，清晰有力。",
        synthesis_user,
    )
    dreamteam_last_synthesis[ALLOWED_USER_ID] = {
        "question": display_question,
        "synthesis": synthesis,
        "responses": all_responses,
    }
    await send_long(update, f"🏆 Dream Team 整合建議：\n\n{synthesis}")
    await update.message.reply_text("💡 用 /dreamteampdf 可以將以上分析生成 PDF 報告")


# ── Dream Team PDF ────────────────────────────────────────────────────────────

async def cmd_dreamteampdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    if not REPORTLAB_AVAILABLE:
        await update.message.reply_text("⚠️ PDF 功能需要 reportlab，請聯絡 VPS 安裝後重啟 bot。")
        return

    data = dreamteam_last_synthesis.get(ALLOWED_USER_ID)
    if not data:
        await update.message.reply_text(
            "⚠️ 未找到 Dream Team 分析結果。\n"
            "請先用 /dreamteam 讓教練分析，再用 /dreamteampdf 生成報告。"
        )
        return

    await update.message.reply_text("📄 生成 Dream Team 策略報告中（約30秒）...")
    loop = asyncio.get_event_loop()

    pdf_prompt = (
        f"根據以下 Dream Team 8位教練嘅分析，生成一份高質素嘅策略報告。\n\n"
        f"議題：{data['question']}\n\n"
        f"各教練意見摘要：\n{data['responses'][:5000]}\n\n"
        f"整合分析：\n{data['synthesis']}\n\n"
        f"輸出格式：繁體中文，第一行係報告標題（唔加 #），章節標題用 # 開頭，"
        f"包含：核心建議（每位教練最重要一點）、共識行動、優先執行順序、風險提示。"
        f"長度 1500-2000字。只輸出報告內容，唔需要其他說明。"
    )

    content = await loop.run_in_executor(
        executor, run_with_system,
        "你係商業策略整合師，輸出高質素 PDF 報告，繁體中文，廣東話語氣。",
        pdf_prompt,
    )
    lines = content.strip().split('\n')
    title = lines[0].lstrip('#').strip() if lines else "Dream Team 策略報告"
    body = '\n'.join(lines[1:]) if len(lines) > 1 else content
    try:
        import io as _io
        pdf_bytes = await loop.run_in_executor(executor, generate_pdf_bytes, title, body)
        safe_name = re.sub(r'[^\w\s]', '', title)[:30].strip().replace(' ', '_') or 'dreamteam_report'
        await update.message.reply_document(
            document=_io.BytesIO(pdf_bytes),
            filename=f"{safe_name}.pdf",
            caption=f"🏆 Dream Team 策略報告\n{title}"
        )
    except Exception as e:
        await update.message.reply_text(f"PDF 生成失敗：{e}\n\n文字版：\n{content[:2000]}")


# ── Report PDF ─────────────────────────────────────────────────────────────────

async def cmd_reportpdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    if not REPORTLAB_AVAILABLE:
        await update.message.reply_text("⚠️ PDF 功能需要 reportlab，請安裝後重啟 bot。")
        return

    report = load_report_buffer()
    items = report.get(str(ALLOWED_USER_ID), [])
    if not items:
        await update.message.reply_text(
            "⚠️ Report 未有內容。\n"
            "請先追問員工並選 [✅ 加入 Report]，再用 /reportpdf 生成 PDF。"
        )
        return

    await update.message.reply_text(f"📄 整合 {len(items)} 條 Q&A，生成 Report PDF 中（約30秒）...")
    loop = asyncio.get_event_loop()

    report_text = ""
    for i, item in enumerate(items, 1):
        agent = item.get('agent', '')
        q = item.get('question', '')
        a = item.get('answer', '')
        report_text += f"【Q{i} 問 {agent}】{q}\n{a}\n\n"

    pdf_prompt = (
        f"以下係 Stanley 同各員工嘅重要追問 Q&A 記錄（共 {len(items)} 條）：\n\n"
        f"{report_text[:5000]}\n\n"
        f"請整理成一份完整嘅工作報告，第一行係報告標題（唔加 #），章節標題用 # 開頭，"
        f"按主題分組，突出重要結論同行動建議。繁體中文，廣東話語氣，長度 1200-2000字。"
        f"只輸出報告內容。"
    )
    content = await loop.run_in_executor(
        executor, run_with_system,
        "你係 Anna，Stanley 團隊文件整理師，幫 Stanley 整理重要 Q&A 成策略報告，廣東話，繁體中文。",
        pdf_prompt,
    )
    lines = content.strip().split('\n')
    title = lines[0].lstrip('#').strip() if lines else "工作報告"
    body = '\n'.join(lines[1:]) if len(lines) > 1 else content
    try:
        import io as _io
        pdf_bytes = await loop.run_in_executor(executor, generate_pdf_bytes, title, body)
        safe_name = re.sub(r'[^\w\s]', '', title)[:30].strip().replace(' ', '_') or 'report'
        await update.message.reply_document(
            document=_io.BytesIO(pdf_bytes),
            filename=f"{safe_name}.pdf",
            caption=f"📄 Report PDF 已生成！共 {len(items)} 條 Q&A\n{title}"
        )
    except Exception as e:
        await update.message.reply_text(f"PDF 生成失敗：{e}\n\n文字版：\n{content[:2000]}")


# ── New output commands ────────────────────────────────────────────────────────

async def cmd_pdf(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    if not REPORTLAB_AVAILABLE:
        await update.message.reply_text(
            "⚠️ PDF 功能需要安裝 reportlab：\n"
            "pip install reportlab\n"
            "然後重啟 bot"
        )
        return
    # PTB v21+ 會將 caption 指令路由到 CommandHandler，需要手動讀取附件
    if update.message and update.message.document:
        doc = update.message.document
        file = await context.bot.get_file(doc.file_id)
        file_bytes = await file.download_as_bytearray()
        loop_pre = asyncio.get_event_loop()
        file_content = await loop_pre.run_in_executor(
            executor, extract_file_content, bytes(file_bytes), doc.mime_type or "", doc.file_name or ""
        )
        if file_content.strip():
            last_content[ALLOWED_USER_ID] = file_content[:8000]
            save_last_content_to_disk(last_content)
    description = " ".join(context.args) if context.args else ""
    error_keywords = ["分析失敗", "無法分析", "讀取失敗", "提取失敗", "❌", "⚠️", "error", "Error"]
    if not description:
        cached = last_content.get(ALLOWED_USER_ID, "")
        # 過濾掉錯誤訊息或無用內容，避免生成垃圾 PDF
        if any(kw in cached for kw in error_keywords):
            cached = ""
        history = conversation_history.get(ALLOWED_USER_ID, [])
        recent_msgs = [
            f"{'Stanley' if m['role']=='user' else 'Amy'}：{m['content'][:800]}"
            for m in history[-4:]
            if not any(kw in m['content'] for kw in error_keywords)
        ]
        recent = "\n".join(recent_msgs)
        if recent:
            # 優先用最近對話，避免舊 last_content 影響內容方向
            description = (
                "根據以下最近對話內容整理成專業 PDF 文件。廣東話，結構清晰。\n\n"
                f"最近對話：\n{recent}\n\n"
            )
            await update.message.reply_text("📄 Anna 根據你嘅內容生成 PDF 中（約30秒）...")
        elif cached:
            description = (
                "根據以下內容整理成專業 PDF 文件。廣東話，結構清晰。\n\n"
                f"原始內容：\n{cached[:3000]}"
            )
            await update.message.reply_text("📄 Anna 根據你嘅內容生成 PDF 中（約30秒）...")
        else:
            await update.message.reply_text(
                "用法：/pdf [描述]\n\n"
                "例子：\n"
                "• /pdf 美容療程服務介紹冊，包括護膚、美白、痛症紓緩\n"
                "• /pdf 痛症管理方案提案，目標客戶30-60歲上班族\n"
                "• /pdf 月度業務報告 2025年6月\n\n"
                "或者先 send 你嘅內容，再 /pdf 即可自動生成。"
            )
            return
    else:
        await update.message.reply_text("📄 Anna 撰寫中（約30秒）...")
    user_msg = (
        f"Stanley 需要一份 PDF 文件。主題：{description}\n\n"
        f"輸出格式規則：\n"
        f"- 第一行係文件標題（唔加 #）\n"
        f"- 章節標題用 # 開頭\n"
        f"- 內容專業完整，繁體中文，廣東話語氣\n"
        f"- 長度：800-1500字\n"
        f"只輸出文件內容，唔需要其他說明。"
    )
    loop = asyncio.get_event_loop()
    content = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Anna'], user_msg)
    lines = content.strip().split('\n')
    title = lines[0].lstrip('#').strip() if lines else description
    body = '\n'.join(lines[1:]) if len(lines) > 1 else content
    try:
        import io as _io
        pdf_bytes = await loop.run_in_executor(executor, generate_pdf_bytes, title, body)
        safe_name = re.sub(r'[^\w\s]', '', title)[:30].strip().replace(' ', '_') or 'document'
        pdf_buf = _io.BytesIO(pdf_bytes)
        await update.message.reply_document(
            document=pdf_buf,
            filename=f"{safe_name}.pdf",
            caption=f"📄 Anna 完成！\n{title}"
        )
    except Exception as e:
        await update.message.reply_text(f"⚠️ PDF 生成失敗：{e}\n\n文字版本：\n{content[:3000]}")


async def cmd_caption(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    description = " ".join(context.args) if context.args else ""
    if not description:
        await update.message.reply_text(
            "用法：/caption [描述]\n\n"
            "例子：\n"
            "• /caption 痛症初診優惠，本月限定，免費諮詢\n"
            "• /caption 美容護膚療程，天然成分，適合敏感肌"
        )
        return
    await update.message.reply_text("🎨 Anna 撰寫中...")
    user_msg = (
        f"幫 Stanley 寫帖文 caption。主題：{description}\n\n"
        f"輸出兩個版本：\n"
        f"【IG 版】emoji + hook + 正文 + CTA + 5-8個hashtag\n"
        f"【小紅書版】吸睛標題 + 生活化正文 + hashtag\n"
        f"直接輸出，唔需要解釋。"
    )
    loop = asyncio.get_event_loop()
    result = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Anna'], user_msg)
    await send_long(update, f"🎨 Anna：\n\n{result}")


async def cmd_reel(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    topic = " ".join(context.args) if context.args else ""
    if not topic:
        await update.message.reply_text(
            "用法：/reel [主題]\n\n"
            "例子：\n"
            "• /reel 三個痛症舒緩小貼士\n"
            "• /reel 美白療程前後對比，真實客戶效果"
        )
        return
    await update.message.reply_text("🎬 Anna 撰寫 Reels 腳本中...")
    user_msg = (
        f"幫 Stanley 寫 60秒 Instagram Reels 腳本。主題：{topic}\n\n"
        f"嚴格按以下格式：\n"
        f"🎣 0-3秒 Hook：[開場白/動作]\n"
        f"💥 3-15秒 衝突/痛點：[內容]\n"
        f"✅ 15-45秒 解方/展示：[內容]\n"
        f"📣 45-60秒 CTA：[行動呼籲]\n\n"
        f"每段包含：① 旁白台詞 ② 畫面描述 ③ 字幕文字\n"
        f"廣東話，生動有力，直接輸出。"
    )
    loop = asyncio.get_event_loop()
    result = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Anna'], user_msg)
    await send_long(update, f"🎬 Anna Reels 腳本：\n\n{result}")


async def cmd_slides(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    if not PPTX_AVAILABLE:
        await update.message.reply_text(
            "⚠️ PowerPoint 功能需要安裝 python-pptx：\n"
            "/root/claude-bot/venv/bin/pip install python-pptx\n"
            "然後重啟 bot"
        )
        return
    # PTB v21+ 會將 caption 指令路由到 CommandHandler，需要手動讀取附件
    if update.message and update.message.document:
        doc = update.message.document
        file = await context.bot.get_file(doc.file_id)
        file_bytes = await file.download_as_bytearray()
        loop_pre = asyncio.get_event_loop()
        file_content = await loop_pre.run_in_executor(
            executor, extract_file_content, bytes(file_bytes), doc.mime_type or "", doc.file_name or ""
        )
        if file_content.strip():
            last_content[ALLOWED_USER_ID] = file_content[:8000]
            save_last_content_to_disk(last_content)
    description = " ".join(context.args) if context.args else ""
    error_keywords = ["分析失敗", "無法分析", "讀取失敗", "提取失敗", "❌", "⚠️", "error", "Error"]
    if not description:
        cached = last_content.get(ALLOWED_USER_ID, "")
        if any(kw in cached for kw in error_keywords):
            cached = ""
        history = conversation_history.get(ALLOWED_USER_ID, [])
        recent_msgs = [
            f"{'Stanley' if m['role']=='user' else 'Amy'}：{m['content'][:800]}"
            for m in history[-4:]
            if not any(kw in m['content'] for kw in error_keywords)
        ]
        recent = "\n".join(recent_msgs)
        if cached or recent:
            description = (
                "根據以下對話內容製作幻燈片。若有優化版本請用優化後嘅版本，廣東話，結構清晰。\n\n"
                + (f"最近對話：\n{recent}\n\n" if recent else "")
                + (f"原始內容：\n{cached[:3000]}" if cached else "")
            )
            await update.message.reply_text("📊 Anna 根據你嘅內容製作幻燈片中（約30秒）...")
        else:
            await update.message.reply_text(
                "用法：/slides [描述]\n\n"
                "例子：\n"
                "• /slides 痛症管理服務介紹，5頁\n"
                "• /slides 美容療程套餐提案，目標客戶30-45歲女性\n"
                "• /slides 月度業務報告 2025年6月\n\n"
                "或者先 send 你嘅內容，再 /slides 即可自動生成。"
            )
            return
    else:
        tool = "Gamma AI" if GAMMA_API_KEY else "Anna (python-pptx)"
        await update.message.reply_text(f"📊 {tool} 製作幻燈片中（約30-60秒）...")

    loop = asyncio.get_event_loop()

    # ── Try Gamma API first ──────────────────────────────────────────────────
    if GAMMA_API_KEY:
        content_msg = (
            f"幫 Stanley 製作一份專業 PowerPoint 簡報。主題：{description}\n\n"
            f"包含：背景分析、核心要點、行動建議。繁體中文，廣東話語氣，5-7頁。"
        )
        content_for_gamma = await loop.run_in_executor(
            executor, run_with_system, AGENT_PROMPTS['Anna'], content_msg, AGENT_MODELS['Anna']
        )
        try:
            pptx_bytes, gamma_url = await loop.run_in_executor(
                executor, generate_gamma_pptx, description, content_for_gamma
            )
            import io as _io
            caption = f"📊 Gamma AI 完成！\n標題：{description}"
            if gamma_url:
                caption += f"\n🔗 Gamma 連結：{gamma_url}"
            if pptx_bytes:
                safe_name = re.sub(r'[^\w\s]', '', description)[:30].strip().replace(' ', '_') or 'slides'
                await update.message.reply_document(
                    document=_io.BytesIO(pptx_bytes),
                    filename=f"{safe_name}.pptx",
                    caption=caption
                )
            else:
                await update.message.reply_text(caption)
            return
        except Exception as e:
            logger.warning(f"Gamma API failed, falling back to python-pptx: {e}")
            await update.message.reply_text("⚠️ Gamma API 暫時唔可用，改用 python-pptx 生成...")

    # ── Fallback: python-pptx ────────────────────────────────────────────────
    user_msg = (
        f"幫 Stanley 製作 PowerPoint 幻燈片。主題：{description}\n\n"
        f"請輸出以下 JSON 格式（只輸出 JSON，唔可以有任何其他文字）：\n"
        f'{{"title":"主標題","subtitle":"副標題或日期","slides":[{{"title":"頁標題","points":["重點一","重點二","重點三"]}},{{"title":"頁標題2","points":["重點一","重點二"]}}]}}\n\n'
        f"要求：共 5-7 頁，每頁 3-5 個重點，繁體中文，廣東話語氣，內容專業落地。"
    )
    raw = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Anna'], user_msg, AGENT_MODELS['Anna'])

    title, subtitle, slides = "", "", []
    try:
        json_match = re.search(r'\{.*\}', raw, re.DOTALL)
        if json_match:
            data = json.loads(json_match.group())
            title = data.get("title", "")
            subtitle = data.get("subtitle", "")
            slides = data.get("slides", [])
    except Exception:
        pass

    if not title or not slides:
        await update.message.reply_text(f"⚠️ 內容解析失敗\n\n原始內容：\n{raw[:2000]}")
        return
    try:
        import io as _io
        pptx_bytes = await loop.run_in_executor(executor, generate_pptx_bytes, title, subtitle, slides)
        safe_name = re.sub(r'[^\w\s]', '', title)[:30].strip().replace(' ', '_') or 'slides'
        await update.message.reply_document(
            document=_io.BytesIO(pptx_bytes),
            filename=f"{safe_name}.pptx",
            caption=f"📊 Anna 完成！{len(slides)} 頁幻燈片\n標題：{title}"
        )
    except Exception as e:
        await update.message.reply_text(f"⚠️ PowerPoint 生成失敗：{e}")


async def cmd_copy(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    request = " ".join(context.args) if context.args else ""
    if not request:
        await update.message.reply_text(
            "用法：/copy [需求]\n\n"
            "例子：\n"
            "• /copy WhatsApp 跟進訊息，客戶上週做完痛症療程\n"
            "• /copy 廣告標題，女性美白療程，A/B 兩版\n"
            "• /copy 電郵邀請客戶回購，語氣溫暖"
        )
        return
    await update.message.reply_text("✍️ Anna 撰寫中...")
    user_msg = (
        f"Stanley 需要以下文案：{request}\n\n"
        f"輸出完整、可直接使用嘅最終版本。廣東話，貼地有力。"
    )
    loop = asyncio.get_event_loop()
    result = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Anna'], user_msg)
    await send_long(update, f"✍️ Anna：\n\n{result}")


# ── Message & file handlers ────────────────────────────────────────────────────

async def send_long(update: Update, text: str):
    for i in range(0, len(text), 4096):
        await update.message.reply_text(text[i: i + 4096])


async def maybe_generate_file(update: Update, agent: str, task: str, loop, extra_context: str = "") -> tuple[str, str] | None:
    """若 Anna 被派去生成文件，直接用 Python 生成器，唔經 CLI 工具。"""
    if agent != "Anna":
        return None
    task_lower = task.lower()

    # Build full_task: if other agents have contributed data, prepend it so Anna can use it
    full_task = task
    if extra_context:
        full_task = (
            f"以下係其他員工已完成嘅研究同分析成果，你必須將呢啲內容整合進你嘅成品，唔好忽略任何重要數據：\n\n"
            f"{extra_context}\n\n{'='*50}\n\n"
            f"原始任務：\n{task}"
        )

    if any(kw in task_lower for kw in ["pdf", "pdf文件", "pdf文檔", "報告", "提案"]):
        if not REPORTLAB_AVAILABLE:
            return None
        await update.message.reply_text("📄 Anna 整合所有員工成果，生成 PDF 中（約60秒）...")
        user_msg = (
            f"Stanley 需要一份高質素嘅 PDF 文件。主題及內容如下：\n{full_task}\n\n"
            f"輸出格式：第一行係文件標題（唔加 #），章節標題用 # 開頭，每個章節有足夠嘅細節同數據支撐，繁體中文，廣東話語氣，長度 1500-2500字。"
            f"要求：內容必須完整詳盡，唔好用省略或者草草帶過；如有數據、洞察、建議行動，全部必須列出。只輸出文件內容，唔需要其他說明。"
        )
        content = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Anna'], user_msg, AGENT_MODELS['Anna'])
        lines = content.strip().split('\n')
        title = lines[0].lstrip('#').strip() if lines else "文件"
        body = '\n'.join(lines[1:]) if len(lines) > 1 else content
        try:
            import io as _io
            pdf_bytes = await loop.run_in_executor(executor, generate_pdf_bytes, title, body)
            safe_name = re.sub(r'[^\w\s]', '', title)[:30].strip().replace(' ', '_') or 'document'
            await update.message.reply_document(
                document=_io.BytesIO(pdf_bytes),
                filename=f"{safe_name}.pdf",
                caption=f"📄 Anna 完成！\n{title}"
            )
            return "Anna", f"PDF 已生成並發送：{title}"
        except Exception as e:
            return "Anna", f"PDF 生成失敗：{e}\n\n文字版：\n{content[:2000]}"

    if any(kw in task_lower for kw in ["slides", "幻燈片", "powerpoint", "簡報", "ppt", ".pptx"]):
        if not PPTX_AVAILABLE and not GAMMA_API_KEY:
            return None
        tool = "Gamma AI" if GAMMA_API_KEY else "Anna"
        await update.message.reply_text(f"📊 {tool} 製作幻燈片中（約30-60秒）...")

        # Try Gamma first
        if GAMMA_API_KEY:
            content_msg = (
                f"幫 Stanley 製作一份專業 PowerPoint 簡報。內容如下：\n{full_task}\n\n"
                f"包含：背景分析、核心要點、行動建議。繁體中文，廣東話語氣，5-7頁。"
            )
            content_for_gamma = await loop.run_in_executor(
                executor, run_with_system, AGENT_PROMPTS['Anna'], content_msg, AGENT_MODELS['Anna']
            )
            try:
                pptx_bytes, gamma_url = await loop.run_in_executor(
                    executor, generate_gamma_pptx, task[:80], content_for_gamma
                )
                import io as _io
                caption = f"📊 Gamma AI 完成！"
                if gamma_url:
                    caption += f"\n🔗 {gamma_url}"
                if pptx_bytes:
                    await update.message.reply_document(
                        document=_io.BytesIO(pptx_bytes),
                        filename="presentation.pptx",
                        caption=caption
                    )
                else:
                    await update.message.reply_text(caption)
                return "Anna", f"Slides 已用 Gamma AI 生成並發送"
            except Exception as e:
                logger.warning(f"Gamma API failed in dispatch: {e}")

        if not PPTX_AVAILABLE:
            return None
        user_msg = (
            f"幫 Stanley 製作 PowerPoint 幻燈片。主題及內容如下：\n{full_task}\n\n"
            f"請輸出以下 JSON 格式（只輸出 JSON，唔可以有任何其他文字）：\n"
            f'{{"title":"主標題","subtitle":"副標題","slides":[{{"title":"頁標題","points":["重點一","重點二","重點三"]}}]}}\n\n'
            f"共 5-7 頁，每頁 3-5 個重點，繁體中文，廣東話語氣。"
        )
        raw = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Anna'], user_msg, AGENT_MODELS['Anna'])
        title, subtitle, slides = "", "", []
        try:
            json_match = re.search(r'\{.*\}', raw, re.DOTALL)
            if json_match:
                data = json.loads(json_match.group())
                title = data.get("title", "")
                subtitle = data.get("subtitle", "")
                slides = data.get("slides", [])
        except Exception:
            pass
        if title and slides:
            try:
                import io as _io
                pptx_bytes = await loop.run_in_executor(executor, generate_pptx_bytes, title, subtitle, slides)
                safe_name = re.sub(r'[^\w\s]', '', title)[:30].strip().replace(' ', '_') or 'slides'
                await update.message.reply_document(
                    document=_io.BytesIO(pptx_bytes),
                    filename=f"{safe_name}.pptx",
                    caption=f"📊 Anna 完成！{len(slides)} 頁幻燈片\n標題：{title}"
                )
                return "Anna", f"Slides 已生成並發送：{title}，共 {len(slides)} 頁"
            except Exception as e:
                return "Anna", f"Slides 生成失敗：{e}"

    if any(kw in task_lower for kw in ["landing page", "landingpage", "html", "網頁"]):
        await update.message.reply_text("🎨 Anna 設計 Landing Page 中（約30秒）...")
        user_msg = (
            f"Stanley 需要一個 Landing Page。內容如下：\n{full_task}\n\n"
            f"生成完整 HTML Landing Page（單一文件），用 Tailwind CSS CDN，"
            f"包含 Hero、服務介紹、痛點解決、Testimonial、CTA，繁體中文，手機優先。只輸出純 HTML。"
        )
        html = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Anna'], user_msg, AGENT_MODELS['Anna'])
        html_match = re.search(r'```(?:html)?\s*(<!DOCTYPE.*?</html>)\s*```', html, re.DOTALL | re.IGNORECASE)
        if html_match:
            html = html_match.group(1)
        import tempfile as _tf, os as _os
        with _tf.NamedTemporaryFile(mode='w', suffix='.html', delete=False, encoding='utf-8') as f:
            f.write(html)
            tmp = f.name
        try:
            with open(tmp, 'rb') as f:
                await update.message.reply_document(document=f, filename='landing_page.html',
                                                     caption="🎨 Anna 完成！用瀏覽器打開預覽。")
        finally:
            _os.unlink(tmp)
        return "Anna", "Landing Page HTML 已生成並發送"

    return None


async def run_with_timeout(tasks: list, update: Update, timeout: int = 240, heartbeat: int = 90):
    """Run agent tasks with heartbeat + timeout. Returns results or None on timeout."""
    async def _heartbeat():
        await asyncio.sleep(heartbeat)
        try:
            await update.message.reply_text("⏳ 仲處理緊，Agent 需要多啲時間，請再等一陣...")
        except Exception:
            pass

    hb = asyncio.create_task(_heartbeat())
    try:
        results = list(await asyncio.wait_for(asyncio.gather(*tasks), timeout=timeout))
        hb.cancel()
        return results
    except asyncio.TimeoutError:
        hb.cancel()
        await update.message.reply_text("⚠️ 任務超時（超過 4 分鐘），請重新發送指令再試。")
        return None
    except Exception as e:
        hb.cancel()
        logger.error(f"Agent task error: {e}")
        await update.message.reply_text("⚠️ 處理出錯，請重試。")
        return None


async def run_with_progress_tracker(
    dispatches: list[dict],
    update: Update,
    loop,
    phase_label: str = "任務進度",
    timeout: int = 240,
) -> list[tuple[str, str]] | None:
    """
    Run agents in parallel. Shows a live ✅ checklist that edits in-place as each completes.
    Posts each agent's result via send_long the moment it finishes.
    Returns list of (agent_name, reply_text) in completion order, or None on timeout.
    """
    if not dispatches:
        return []

    agent_names = [d["agent"] for d in dispatches]
    status: dict[str, str] = {name: "⬜" for name in agent_names}
    completed: list[tuple[str, str]] = []

    def format_tracker() -> str:
        lines = [f"📋 {phase_label}："]
        for name in agent_names:
            emoji_a = AGENT_EMOJI.get(name, "🤖")
            lines.append(f"  {status[name]} {emoji_a} {name}")
        return "\n".join(lines)

    try:
        tracker_msg = await update.message.reply_text(format_tracker())
    except Exception:
        tracker_msg = None

    async def run_one(agent_name: str, task: str) -> tuple[str, str]:
        try:
            return await loop.run_in_executor(executor, agent_call, agent_name, task)
        except Exception as e:
            logger.error(f"Agent {agent_name} error: {e}")
            return (agent_name, f"[ERROR] {str(e)}")

    coros = [run_one(d["agent"], d["task"]) for d in dispatches]

    async def run_all_tracked():
        for fut in asyncio.as_completed(coros):
            try:
                agent_name, reply_text = await fut
            except Exception as e:
                logger.error(f"Tracker fut error: {e}")
                continue
            status[agent_name] = "✅"
            if tracker_msg:
                try:
                    await tracker_msg.edit_text(format_tracker())
                except Exception:
                    pass
            completed.append((agent_name, reply_text))
            if "[QUOTA_EXCEEDED]" not in reply_text and "[ERROR]" not in reply_text:
                emoji = AGENT_EMOJI.get(agent_name, "🤖")
                await send_long(update, f"{emoji} {agent_name}：\n{reply_text}")

    try:
        await asyncio.wait_for(run_all_tracked(), timeout=timeout)
    except asyncio.TimeoutError:
        await update.message.reply_text("⚠️ 任務超時（超過 4 分鐘），請重新發送指令再試。")
        return None

    return completed


async def dispatch_with_content(update: Update, file_content: str, user_intent: str):
    loop = asyncio.get_event_loop()
    history = conversation_history.setdefault(ALLOWED_USER_ID, [])

    last_content[ALLOWED_USER_ID] = file_content[:8000]
    save_last_content_to_disk(last_content)

    combined_input = (
        f"Stanley 上傳咗一份文件，內容如下：\n\n"
        f"{'='*40}\n{file_content[:6000]}\n{'='*40}\n\n"
        f"Stanley 嘅指示：{user_intent or '請分析呢份文件，並產出對我業務最有用嘅成品。'}"
    )

    raw = await loop.run_in_executor(
        executor, run_with_system, AMY_DISPATCH_SYSTEM,
        f"Stanley 最新指令：{combined_input}",
        MODEL_FAST
    )

    if "[QUOTA_EXCEEDED]" in raw:
        await update.message.reply_text(QUOTA_EXCEEDED_MSG)
        return

    try:
        match = re.search(r'\{.*\}', raw, re.DOTALL)
        data = json.loads(match.group()) if match else {}
    except Exception:
        await send_long(update, f"{AGENT_EMOJI['Amy']} Amy：\n{raw}")
        return

    if data.get("direct_reply"):
        await send_long(update, f"{AGENT_EMOJI['Amy']} Amy：{data['direct_reply']}")
        return

    dispatches = [d for d in data.get("dispatch", []) if d.get("agent") in AGENT_PROMPTS]
    amy_msg = data.get("amy_message", "")

    if not dispatches:
        await send_long(update, f"{AGENT_EMOJI['Amy']} Amy：{amy_msg or raw}")
        return

    agents_names = "、".join(d["agent"] for d in dispatches)
    await send_long(update, f"{AGENT_EMOJI['Amy']} Amy：{amy_msg}\n（調度中：{agents_names}）")

    full_results = ""

    # Step 1: Separate file agents (Anna PDF/slides/html) from analysis agents
    file_dispatches = []
    normal_dispatches = []
    _file_keywords = ["pdf", "pdf文件", "pdf文檔", "slides", "幻燈片", "powerpoint", "簡報", "ppt", ".pptx", "landing page", "landingpage", "html", "網頁", "報告", "提案"]
    for d in dispatches:
        is_file = d["agent"] == "Anna" and any(kw in d["task"].lower() for kw in _file_keywords)
        (file_dispatches if is_file else normal_dispatches).append(d)

    # Step 2: Run analysis agents in parallel first — collect their results
    if normal_dispatches:
        results = await run_with_progress_tracker(normal_dispatches, update, loop, phase_label="員工任務進度")
        if results is None:
            return
        for agent_name, reply_text in results:
            if "[QUOTA_EXCEEDED]" in reply_text:
                await update.message.reply_text(QUOTA_EXCEEDED_MSG)
                return
            emoji = AGENT_EMOJI.get(agent_name, "🤖")
            full_results += f"{emoji} {agent_name}：\n{reply_text}\n\n"

    # Step 3: Run file agents with ALL analysis results so Anna has full context
    combined_context = full_results  # includes all analysis agents' output

    # Amy organises all collected data into a structured creative brief for Anna
    if file_dispatches and combined_context:
        await update.message.reply_text("📋 Amy 整理所有員工成果，為 Anna 準備完整製作簡報...")
        brief_prompt = (
            f"你係 Amy，需要整理一份完整嘅製作簡報俾 Anna。\n\n"
            f"Stanley 嘅原始指令：{combined_input[:1500]}\n\n"
            f"各員工已完成嘅成果：\n{combined_context[:5000]}\n\n"
            f"請整理成一份結構清晰、內容完整嘅製作簡報，包含：\n"
            f"① 成品目標同定位（Stanley 想要咩、用於咩場合）\n"
            f"② 核心數據同洞察（Leo 嘅競品分析、市場數據，全部列出）\n"
            f"③ 策略重點（Small 嘅策略建議，具體列出）\n"
            f"④ 客戶話術同轉化點（Tony 嘅建議，具體列出）\n"
            f"⑤ 建議內容結構（Anna 應該分幾個章節、每個章節講咩）\n"
            f"⑥ 語氣同風格要求\n\n"
            f"輸出一份清晰、完整嘅簡報，讓 Anna 可以完全根據呢份簡報製作最終成品，唔需要再猜或補充。"
        )
        _, structured_brief = await loop.run_in_executor(
            executor, run_with_system, AGENT_PROMPTS['Amy'], brief_prompt, AGENT_MODELS['Amy']
        )
        if "[QUOTA_EXCEEDED]" not in structured_brief:
            combined_context = f"【Amy 製作簡報】\n{structured_brief}"

    for d in file_dispatches:
        file_result = await maybe_generate_file(update, d["agent"], d["task"], loop, extra_context=combined_context)
        if file_result:
            agent_name, reply_text = file_result
            full_results += f"{AGENT_EMOJI.get(agent_name,'🤖')} {agent_name}：\n{reply_text}\n\n"
        else:
            enriched_task = (f"{combined_context}\n\n原始任務：\n{d['task']}") if combined_context else d["task"]
            _, reply_text = await loop.run_in_executor(executor, agent_call, d["agent"], enriched_task)
            emoji = AGENT_EMOJI.get(d["agent"], "🤖")
            msg = f"{emoji} {d['agent']}：\n{reply_text}"
            full_results += msg + "\n\n"
            await send_long(update, msg)

    consolidate_user = (
        f"Stanley 嘅原始指令：{combined_input[:2000]}\n\n"
        f"各 Agent 回覆：\n{full_results}\n\n"
        f"請整合所有 Agent 嘅回覆，突出最重要嘅成品同行動點。"
    )
    summary = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Amy'], consolidate_user, AGENT_MODELS['Amy'])
    if "[QUOTA_EXCEEDED]" in summary:
        await update.message.reply_text(QUOTA_EXCEEDED_MSG)
        return
    await send_long(update, f"{AGENT_EMOJI['Amy']} Amy 整合：\n{summary}")

    # Save agent outputs to last_content so follow-up "optimize/redo" commands have context
    last_content[ALLOWED_USER_ID] = full_results[:8000]
    save_last_content_to_disk(last_content)

    history.append({"role": "user", "content": combined_input[:3000]})
    history.append({"role": "assistant", "content": summary})
    save_history(conversation_history)

    warning = get_usage_warning()
    if warning:
        await update.message.reply_text(warning)


async def handle_photo(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return

    photo = update.message.photo[-1]
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    await update.message.reply_text("🖼️ 收到圖片，分析中...")

    file = await context.bot.get_file(photo.file_id)
    img_bytes = await file.download_as_bytearray()
    user_intent = update.message.caption or ""

    loop = asyncio.get_event_loop()
    file_content = await loop.run_in_executor(executor, extract_file_content, bytes(img_bytes), "image/jpeg", "photo.jpg")

    await dispatch_with_content(update, file_content, user_intent)


async def handle_document(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    doc = update.message.document
    if not doc:
        return

    mime = doc.mime_type or ""
    filename = doc.file_name or ""

    type_labels = {
        "application/pdf": "📄 PDF",
        "image/jpeg": "🖼️ 圖片",
        "image/png": "🖼️ 圖片",
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document": "📝 Word",
        "text/plain": "📝 文字",
        "text/csv": "📊 CSV",
    }
    label = type_labels.get(mime, f"📎 {filename}")
    await update.message.reply_text(f"{label} 收到，讀取中...")
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")

    file = await context.bot.get_file(doc.file_id)
    file_bytes = await file.download_as_bytearray()
    user_intent = update.message.caption or ""

    supported = (
        mime == "application/pdf" or
        mime.startswith("image/") or
        mime in ("text/plain", "text/csv", "text/markdown", "application/json",
                 "application/vnd.openxmlformats-officedocument.wordprocessingml.document") or
        any(filename.endswith(ext) for ext in (".txt", ".csv", ".md", ".json", ".docx"))
    )

    if not supported:
        await update.message.reply_text(
            f"文件類型 {mime or filename} 暫時唔支援直接讀取。\n"
            "支援：PDF、圖片、Word(.docx)、TXT、CSV、JSON、Markdown\n\n"
            "可以將內容複製貼上俾我。"
        )
        return

    loop = asyncio.get_event_loop()
    file_content = await loop.run_in_executor(executor, extract_file_content, bytes(file_bytes), mime, filename)

    if not file_content.strip():
        await update.message.reply_text("讀取文件失敗，請試下複製內容貼俾我。")
        return

    await dispatch_with_content(update, file_content, user_intent)


async def handle_voice(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    voice = update.message.voice or update.message.audio
    if not voice:
        return
    if not FASTER_WHISPER_AVAILABLE and not (OPENAI_AVAILABLE and OPENAI_API_KEY):
        await update.message.reply_text(
            "⚠️ 語音功能未啟用。\n\n"
            "VPS 安裝免費語音轉錄：\n"
            "pip install faster-whisper\n"
            "systemctl restart claude-bot"
        )
        return
    await update.message.reply_text("🎙️ 收到錄音，轉錄中...")
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    file = await context.bot.get_file(voice.file_id)
    audio_bytes = await file.download_as_bytearray()
    loop = asyncio.get_event_loop()
    transcript = await loop.run_in_executor(
        executor, transcribe_audio, bytes(audio_bytes), "voice.ogg"
    )
    if transcript.startswith("["):
        await update.message.reply_text(f"⚠️ {transcript}")
        return
    await update.message.reply_text(f"🎙️ 轉錄結果：\n\n{transcript}")
    await dispatch_with_content(
        update,
        f"[語音訊息轉錄內容]\n{transcript}",
        "請根據以上語音內容，理解 Stanley 嘅意圖，產出相應成品或回應。"
    )


def detect_agent_from_message(text: str) -> str:
    """Parse emoji in first 60 chars to identify which agent sent a message."""
    if not text:
        return ""
    preview = text[:60]
    for name, emoji in AGENT_EMOJI.items():
        if emoji in preview:
            return name
    return ""


async def handle_reply_to_agent(update: Update, context: ContextTypes.DEFAULT_TYPE, agent_name: str, original_text: str):
    """User swiped left and replied to an agent message — route follow-up directly to that agent."""
    user_question = update.message.text
    emoji = AGENT_EMOJI.get(agent_name, "🤖")
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    all_outputs = agent_outputs.get(ALLOWED_USER_ID, {})
    agent_ctx = all_outputs.get(agent_name, original_text[:3000])
    # 跨員工 context：把其他員工嘅成果一起帶入
    other_ctx = ""
    for other_name, other_output in all_outputs.items():
        if other_name != agent_name:
            other_emoji = AGENT_EMOJI.get(other_name, "🤖")
            other_ctx += f"\n\n【{other_emoji} {other_name} 嘅成果】\n{other_output[:1500]}"
    prompt = (
        f"以下係你之前完成嘅分析成果：\n\n{agent_ctx}\n\n"
        + (f"{'='*30}\n【其他員工嘅相關成果（供參考）】{other_ctx}\n\n{'='*30}\n\n" if other_ctx else "")
        + f"Stanley 追問：{user_question}\n\n"
        f"請根據你上面嘅分析（以及其他員工嘅成果），詳細回答 Stanley 嘅追問。如有需要可跨員工洞見提供更全面答案。"
    )
    loop = asyncio.get_event_loop()
    _, reply_text = await loop.run_in_executor(executor, agent_call, agent_name, prompt)
    await send_long(update, f"{emoji} {agent_name}：\n{reply_text}")
    pending_report_entry[ALLOWED_USER_ID] = {
        "agent": agent_name,
        "question": user_question,
        "answer": reply_text,
    }
    buttons = [[
        InlineKeyboardButton("✅ 加入 Report", callback_data="report_action:add"),
        InlineKeyboardButton("❌ 唔使", callback_data="report_action:skip"),
    ]]
    await update.message.reply_text("要加入 Report 嗎？", reply_markup=InlineKeyboardMarkup(buttons))


async def handle_report_action_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    if query.from_user.id != ALLOWED_USER_ID:
        return
    action = query.data.split(":", 1)[1]
    entry = pending_report_entry.pop(ALLOWED_USER_ID, None)
    if action == "add" and entry:
        import datetime
        entry["ts"] = datetime.datetime.now().strftime("%Y-%m-%d %H:%M")
        key = str(ALLOWED_USER_ID)
        report_buffer.setdefault(key, []).append(entry)
        save_report_buffer(report_buffer)
        await query.edit_message_text("✅ 已加入 Report！用 /myreport 睇累積內容。")
    else:
        await query.edit_message_text("👌 OK，純粹問問。")


async def cmd_myreport(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    key = str(ALLOWED_USER_ID)
    entries = report_buffer.get(key, [])
    if not entries:
        await update.message.reply_text("📋 Report 係空嘅，仲未加入任何內容。\n\n追問員工後，tap ✅ 加入 Report 先。")
        return
    lines = ["📋 *累積 Report*\n"]
    for i, e in enumerate(entries, 1):
        emoji = AGENT_EMOJI.get(e["agent"], "🤖")
        lines.append(f"*{i}. {emoji} {e['agent']}* — {e.get('ts','')}")
        lines.append(f"❓ {e['question']}")
        lines.append(f"💬 {e['answer'][:800]}")
        lines.append("")
    await send_long(update, "\n".join(lines))


async def handle_followup_callback(update: Update, context: ContextTypes.DEFAULT_TYPE):
    query = update.callback_query
    await query.answer()
    if query.from_user.id != ALLOWED_USER_ID:
        return
    agent_name = query.data.split(":", 1)[1]
    outputs = agent_outputs.get(ALLOWED_USER_ID, {})
    agent_ctx = outputs.get(agent_name, last_content.get(ALLOWED_USER_ID, ""))
    follow_up_state[ALLOWED_USER_ID] = {"agent": agent_name, "context": agent_ctx[:4000], "active": True}
    emoji = AGENT_EMOJI.get(agent_name, "🤖")
    await query.message.reply_text(
        f"{emoji} {agent_name} 追問模式啟動！\n\n"
        f"直接問 {agent_name} 任何問題，佢會根據剛才嘅成果回答你。\n"
        f"你問嘅嘢同答案會自動加落去，之後 /pdf /slides /report 生成嘅成品會包含呢啲細節。\n\n"
        f"輸入 /done 退出追問模式。"
    )


async def handle_followup_message(update: Update, context: ContextTypes.DEFAULT_TYPE, state: dict):
    agent_name = state["agent"]
    agent_ctx = state["context"]
    user_question = update.message.text
    loop = asyncio.get_event_loop()
    emoji = AGENT_EMOJI.get(agent_name, "🤖")
    await context.bot.send_chat_action(chat_id=update.effective_chat.id, action="typing")
    all_outputs = agent_outputs.get(ALLOWED_USER_ID, {})
    other_ctx = ""
    for other_name, other_output in all_outputs.items():
        if other_name != agent_name:
            other_emoji = AGENT_EMOJI.get(other_name, "🤖")
            other_ctx += f"\n\n【{other_emoji} {other_name} 嘅成果】\n{other_output[:1200]}"
    prompt = (
        f"以下係你之前完成嘅分析成果：\n\n{agent_ctx}\n\n"
        + (f"{'='*30}\n【其他員工嘅相關成果（供參考）】{other_ctx}\n\n{'='*30}\n\n" if other_ctx else "")
        + f"Stanley 追問：{user_question}\n\n"
        f"請根據你上面嘅分析（以及其他員工嘅成果），詳細回答 Stanley 嘅追問。如有需要可跨員工洞見提供更全面答案。"
    )
    _, reply_text = await loop.run_in_executor(executor, agent_call, agent_name, prompt)
    h = conversation_history.setdefault(ALLOWED_USER_ID, [])
    h.append({"role": "user", "content": f"[追問 {agent_name}] {user_question}"})
    h.append({"role": "assistant", "content": f"{agent_name}：{reply_text}"})
    save_history(conversation_history)
    updated_ctx = f"{agent_ctx}\n\n【追問 {agent_name}】\n問：{user_question}\n答：{reply_text[:3000]}"
    last_content[ALLOWED_USER_ID] = updated_ctx[:8000]
    save_last_content_to_disk(last_content)
    follow_up_state[ALLOWED_USER_ID]["context"] = updated_ctx[:4000]
    await send_long(update, f"{emoji} {agent_name}：\n{reply_text}")
    pending_report_entry[ALLOWED_USER_ID] = {
        "agent": agent_name,
        "question": user_question,
        "answer": reply_text,
    }
    buttons = [[
        InlineKeyboardButton("✅ 加入 Report", callback_data="report_action:add"),
        InlineKeyboardButton("❌ 唔使", callback_data="report_action:skip"),
    ]]
    await update.message.reply_text("要加入 Report 嗎？", reply_markup=InlineKeyboardMarkup(buttons))


async def cmd_done(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    follow_up_state[ALLOWED_USER_ID] = {"active": False}
    await update.message.reply_text(
        "✅ 追問模式結束。\n\n"
        "你嘅追問同答案已經加入記憶。\n"
        "依家可以用 /pdf、/slides、/report 等指令生成包含最新細節嘅成品。"
    )


async def handle_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    user_message = update.message.text
    if not user_message:
        return

    # 向左掃 reply 員工訊息 → 直接追問該員工
    replied = update.message.reply_to_message
    if replied and replied.from_user and replied.from_user.is_bot:
        agent_name = detect_agent_from_message(replied.text or "")
        if agent_name:
            await handle_reply_to_agent(update, context, agent_name, replied.text or "")
            return

    # Follow-up mode: route to specific agent Q&A
    state = follow_up_state.get(ALLOWED_USER_ID, {})
    if state.get("active"):
        if user_message.startswith("/"):
            follow_up_state[ALLOWED_USER_ID] = {"active": False}
        else:
            await handle_followup_message(update, context, state)
            return

    chat_id = update.effective_chat.id
    await context.bot.send_chat_action(chat_id=chat_id, action="typing")

    history = conversation_history.setdefault(ALLOWED_USER_ID, [])
    loop = asyncio.get_event_loop()

    # 長文字（>200字）自動記住，持久化到磁碟
    if len(user_message) > 200:
        last_content[ALLOWED_USER_ID] = user_message
        save_last_content_to_disk(last_content)

    history_text = ""
    for msg in history[-MAX_HISTORY:]:
        role = "Stanley" if msg["role"] == "user" else "Amy"
        history_text += f"{role}: {msg['content']}\n\n"

    # 短跟進指令自動附上上次內容，解決 Amy「唔記得」問題
    cached = last_content.get(ALLOWED_USER_ID, "")
    content_hint = (
        f"\n\n[Stanley 上次分享嘅內容（可直接使用）：\n{cached[:3000]}]\n"
        if cached and len(user_message) < 200 else ""
    )

    dispatch_user = (
        f"{'對話歷史：' + chr(10) + history_text if history_text else ''}"
        f"{content_hint}"
        f"Stanley 最新指令：{user_message}"
    )
    raw = await loop.run_in_executor(executor, run_with_system, AMY_DISPATCH_SYSTEM, dispatch_user, MODEL_FAST)

    if "[QUOTA_EXCEEDED]" in raw:
        await update.message.reply_text(QUOTA_EXCEEDED_MSG)
        return

    try:
        match = re.search(r'\{.*\}', raw, re.DOTALL)
        data = json.loads(match.group()) if match else {}
    except Exception:
        history.append({"role": "user", "content": user_message})
        history.append({"role": "assistant", "content": raw})
        save_history(conversation_history)
        await send_long(update, f"{AGENT_EMOJI['Amy']} Amy：{raw}")
        return

    if data.get("direct_reply"):
        reply = f"{AGENT_EMOJI['Amy']} Amy：{data['direct_reply']}"
        history.append({"role": "user", "content": user_message})
        history.append({"role": "assistant", "content": reply})
        save_history(conversation_history)
        await send_long(update, reply)
        return

    dispatches = [d for d in data.get("dispatch", []) if d.get("agent") in AGENT_PROMPTS]
    actions = [a for a in data.get("actions", []) if a.get("type") and a.get("param")]
    amy_msg = data.get("amy_message", "")

    if not dispatches and not actions:
        reply = f"{AGENT_EMOJI['Amy']} Amy：{amy_msg or raw}"
        history.append({"role": "user", "content": user_message})
        history.append({"role": "assistant", "content": reply})
        save_history(conversation_history)
        await send_long(update, reply)
        return

    announce_parts = [f"{AGENT_EMOJI['Amy']} Amy：{amy_msg}"] if amy_msg else []
    if actions:
        action_labels = {"scrape_ig": "📡 IG抓取", "scrape_xhs": "📕 小紅書", "scrape_web": "🌐 網頁", "scrape_news": "🔍 新聞搜尋"}
        announce_parts.append("（執行中：" + "、".join(action_labels.get(a["type"], a["type"]) for a in actions) + "）")
    if dispatches:
        announce_parts.append("（調度中：" + "、".join(d["agent"] for d in dispatches) + "）")
    await send_long(update, "\n".join(announce_parts))

    # ── Handle real-execution actions ──────────────────────────────────────────
    apify_ok = APIFY_AVAILABLE and APIFY_TOKEN
    action_context = ""  # accumulates scrape + analysis results to pass to file agents later

    for action in actions:
        atype = action["type"]
        param = action["param"].strip().lstrip("@")
        # scrape_ig uses instaloader (no Apify needed); others need Apify
        if atype in ("scrape_threads", "scrape_xhs", "scrape_web", "scrape_news") and not apify_ok:
            await update.message.reply_text("⚠️ Apify 未設定，請先設定 APIFY_API_TOKEN（見 /status）")
            continue
        def _apify_failed(r: str) -> bool:
            return any(kw in r for kw in ["失敗", "未設定", "不到資料", "Error", "error", "⚠️", "FAILED"])

        if atype == "scrape_fb":
            await update.message.reply_text(f"📘 抓取 Facebook @{param} 中...")
            result = await loop.run_in_executor(executor, scrape_facebook, param)
            await send_long(update, result)
            leo_task = (
                f"以下係競品 Facebook 專頁 @{param} 嘅最新帖文，請做競品分析：\n\n{result}\n\n"
                f"輸出：① 專頁內容定位 ② 爆款帖文規律 ③ 對 Stanley 業務嘅具體啟示 ④ 可以借鑒嘅策略"
            )
            _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
            await send_with_avatar(update, "Leo", analysis)
            action_context += f"【Facebook @{param} 競品分析 by Leo】\n{analysis[:3000]}\n\n"
        elif atype == "scrape_threads":
            await update.message.reply_text(f"🧵 抓取 Threads @{param} 中...")
            result = await loop.run_in_executor(executor, scrape_threads, param)
            await send_long(update, result)
            leo_task = (
                f"以下係競品 Threads 帳號 @{param} 嘅最新內容，請做競品分析：\n\n{result}\n\n"
                f"輸出：① 帳號內容定位 ② 爆款帖文規律 ③ 對 Stanley 業務嘅具體啟示 ④ 可以借鑒嘅策略"
            )
            _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
            await send_with_avatar(update, "Leo", analysis)
            action_context += f"【Threads @{param} 競品分析 by Leo】\n{analysis[:3000]}\n\n"
        elif atype == "scrape_ig":
            await update.message.reply_text(f"📡 抓取 IG @{param} 中（約30-60秒）...")
            result = await loop.run_in_executor(executor, scrape_instagram, param)
            await send_long(update, result)
            if _apify_failed(result):
                await update.message.reply_text(
                    "❌ IG 數據抓取失敗，Leo 分析已取消（節省 tokens）。\n"
                    "請去 apify.com → Billing 確認 credits 充足後再試。"
                )
                continue
            leo_task = (
                f"以下係競品 IG 帳號 @{param} 嘅最新帖文數據，請做完整競品分析：\n\n{result}\n\n"
                f"輸出：① 帳號整體定位 ② 爆款帖文規律 ③ 對 Stanley 美容/痛症業務嘅具體啟示 ④ 可以抄嘅策略"
            )
            _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
            await send_with_avatar(update, "Leo", analysis)
            action_context += f"【IG @{param} 原始數據】\n{result[:2000]}\n\n【Leo 競品分析】\n{analysis[:3000]}\n\n"
        elif atype == "scrape_xhs":
            await update.message.reply_text(f"📕 抓取小紅書「{param}」中（約30-60秒）...")
            result = await loop.run_in_executor(executor, scrape_xhs, param)
            await send_long(update, result)
            anna_task = (
                f"以下係小紅書「{param}」帖文內容，請分析趨勢：\n\n{result[:3000]}\n\n"
                f"輸出：① 爆款標題規律 ② 熱門話題角度 ③ Stanley 可以用嘅3個文案方向"
            )
            _, analysis = await loop.run_in_executor(executor, agent_call, "Anna", anna_task)
            await send_with_avatar(update, "Anna", analysis)
            action_context += f"【小紅書「{param}」趨勢分析 by Anna】\n{analysis[:3000]}\n\n"
        elif atype == "scrape_web":
            await update.message.reply_text(f"🌐 抓取 {param} 中...")
            result = await loop.run_in_executor(executor, scrape_webpage, param)
            await send_long(update, result)
            action_context += f"【網頁內容 {param}】\n{result[:3000]}\n\n"
        elif atype == "scrape_news":
            await update.message.reply_text(f"🔍 搜尋「{param}」新聞中...")
            result = await loop.run_in_executor(executor, scrape_google, param)
            await send_long(update, result)
            leo_task = f"以下係「{param}」嘅最新搜尋結果，請分析對 Stanley 業務嘅影響同機會：\n\n{result[:3000]}"
            _, analysis = await loop.run_in_executor(executor, agent_call, "Leo", leo_task)
            await send_with_avatar(update, "Leo", analysis)
            action_context += f"【新聞「{param}」分析 by Leo】\n{analysis[:3000]}\n\n"
        elif atype == "product_research":
            topic = param
            await update.message.reply_text(f"🔬 開始市場研究：「{topic}」\n正在擷取真實數據...")
            # Always use web_search (DuckDuckGo) as primary — reliable, free, no Apify needed
            search_fns = [
                loop.run_in_executor(executor, web_search, topic + " 香港市場趨勢 2025", 6),
                loop.run_in_executor(executor, web_search, topic + " 競品 用戶痛點 新產品 香港", 5),
            ]
            search_results = await asyncio.gather(*search_fns)
            combined_data = "\n\n---\n\n".join(str(r) for r in search_results)[:4000]
            # Supplement with XHS data if Apify available
            if apify_ok:
                xhs_raw = await loop.run_in_executor(executor, scrape_xhs, topic)
                if "失敗" not in xhs_raw and "未設定" not in xhs_raw and "未找到" not in xhs_raw:
                    combined_data += f"\n\n---\n\n【小紅書真實帖文】\n{xhs_raw[:2000]}"
            await update.message.reply_text("✅ 數據擷取完成，Leo + Small 分析緊...")
            leo_task = (
                f"根據以下真實市場數據，分析「{topic}」嘅市場現況：\n\n{combined_data}\n\n"
                f"直接輸出：\n① 市場趨勢同規模\n② 主流產品/服務係乜\n③ 用戶最大痛點\n④ 市場空白位"
            )
            small_task = (
                f"根據以下「{topic}」真實市場數據，識別產品機會：\n\n{combined_data}\n\n"
                f"Stanley 業務背景：香港美容（護膚療程）+ 痛症管理，主要客戶 18-65歲。\n\n"
                f"直接輸出（唔係計劃，係具體可執行嘅結果）：\n"
                f"① 3個市場機會（每個2句描述）\n"
                f"② 2個新產品/服務建議（名稱 + 目標客群 + 核心賣點 + 建議定價）\n"
                f"③ 本週可以立即開始做嘅第一步"
            )
            (_, market_analysis), (_, product_ideas) = await asyncio.gather(
                loop.run_in_executor(executor, agent_call, "Leo", leo_task),
                loop.run_in_executor(executor, agent_call, "Small", small_task),
            )
            await send_long(update, f"📊 Leo — 市場分析：\n{market_analysis}")
            await send_with_avatar(update, "Small", product_ideas)
            action_context += f"【市場研究「{topic}」— Leo 市場分析】\n{market_analysis[:2000]}\n\n【Small 產品機會】\n{product_ideas[:2000]}\n\n"

    full_results = ""

    # ── 分三個 Phase 執行 ──────────────────────────────────────────────────────
    _file_keywords = ["pdf", "pdf文件", "pdf文檔", "slides", "幻燈片", "powerpoint", "簡報", "ppt", ".pptx", "landing page", "landingpage", "html", "網頁", "報告", "提案"]
    file_dispatches = []
    research_dispatches = []
    strategy_dispatches = []
    for d in dispatches:
        is_file = d["agent"] == "Anna" and any(kw in d["task"].lower() for kw in _file_keywords)
        if is_file:
            file_dispatches.append(d)
        elif d["agent"] in RESEARCH_AGENTS:
            research_dispatches.append(d)
        else:
            strategy_dispatches.append(d)

    research_results = ""   # Phase 1 成果，傳俾 Phase 2
    all_agents_ran = []

    # ── Phase 1：研究員（Leo / Kai）——搜集數據 ──────────────────────────────
    if research_dispatches:
        r_names = "、".join(d["agent"] for d in research_dispatches)
        await send_long(update, f"👩‍💼 Amy：【Phase 1】派出研究員 {r_names}，搜集數據中...")
        results = await run_with_progress_tracker(research_dispatches, update, loop, phase_label="Phase 1 研究進度")
        if results is None:
            return
        for agent_name, reply_text in results:
            if "[QUOTA_EXCEEDED]" in reply_text:
                await update.message.reply_text(QUOTA_EXCEEDED_MSG)
                return
            emoji = AGENT_EMOJI.get(agent_name, "🤖")
            full_results += f"{emoji} {agent_name}：\n{reply_text}\n\n"
            agent_outputs.setdefault(ALLOWED_USER_ID, {})[agent_name] = reply_text
            save_agent_outputs_to_disk(agent_outputs)
            research_results += f"【{emoji} {agent_name} 研究成果】\n{reply_text[:2500]}\n\n"
            all_agents_ran.append(agent_name)

    # ── Phase 2：策略師（Small / Tony / Rex / Mia / Toxic）——整合研究出策略 ──
    if strategy_dispatches:
        s_names = "、".join(d["agent"] for d in strategy_dispatches)
        if research_results:
            r_done = "、".join(d["agent"] for d in research_dispatches)
            await send_long(update, f"👩‍💼 Amy：{r_done} 研究完成！將數據傳俾 {s_names}，制定策略中...")
        else:
            await send_long(update, f"👩‍💼 Amy：【Phase 2】派出 {s_names} 制定策略中...")

        augmented = []
        for d in strategy_dispatches:
            task = d["task"]
            if research_results:
                task = (
                    f"【研究員已完成數據搜集，你必須根據以下真實數據制定策略，唔好忽略任何重要發現】\n\n"
                    f"{research_results[:3000]}\n"
                    f"{'='*40}\n\n"
                    f"你的任務：\n{task}"
                )
            augmented.append({"agent": d["agent"], "task": task})

        results = await run_with_progress_tracker(augmented, update, loop, phase_label="Phase 2 策略進度")
        if results is None:
            return
        for agent_name, reply_text in results:
            if "[QUOTA_EXCEEDED]" in reply_text:
                await update.message.reply_text(QUOTA_EXCEEDED_MSG)
                return
            emoji = AGENT_EMOJI.get(agent_name, "🤖")
            full_results += f"{emoji} {agent_name}：\n{reply_text}\n\n"
            agent_outputs.setdefault(ALLOWED_USER_ID, {})[agent_name] = reply_text
            save_agent_outputs_to_disk(agent_outputs)
            all_agents_ran.append(agent_name)

    # Phase 3（Anna）係後面 file_dispatches 處理，唔係喺呢度

    # 追問按鈕（所有 Phase 1 + 2 嘅員工）
    if all_agents_ran:
        buttons = [[InlineKeyboardButton(f"🔍 追問 {n}", callback_data=f"followup:{n}")] for n in all_agents_ran]
        await update.message.reply_text("💬 想追問任何員工？", reply_markup=InlineKeyboardMarkup(buttons))

    # 整合所有成果 → Phase 3 Anna 用
    combined_context = ""
    if action_context:
        combined_context += f"【數據擷取結果】\n{action_context}\n\n"
    if full_results:
        combined_context += f"【Phase 1+2 員工成果】\n{full_results}\n\n"

    # ── Phase 3：Anna 製作成品（收到所有 Phase 1+2 成果後才出手）────────────
    if file_dispatches:
        if combined_context:
            await send_long(update, "👩‍💼 Amy：【Phase 3】所有員工完成！整理成果，交俾 🎨 Anna 製作最終成品中...")
        else:
            await send_long(update, "👩‍💼 Amy：【Phase 3】派出 🎨 Anna 製作成品中...")

    if file_dispatches and combined_context:
        await update.message.reply_text("📋 Amy 整合所有員工成果，為 Anna 準備完整製作簡報...")
        brief_prompt = (
            f"你係 Amy，需要整理一份完整嘅製作簡報俾 Anna。\n\n"
            f"Stanley 嘅原始指令：{user_message}\n\n"
            f"各員工已完成嘅成果：\n{combined_context[:5000]}\n\n"
            f"請整理成一份結構清晰、內容完整嘅製作簡報，包含：\n"
            f"① 成品目標同定位（Stanley 想要咩、用於咩場合）\n"
            f"② 核心數據同洞察（Leo 嘅競品分析、市場數據，全部列出）\n"
            f"③ 策略重點（Small 嘅策略建議，具體列出）\n"
            f"④ 客戶話術同轉化點（Tony 嘅建議，具體列出）\n"
            f"⑤ 建議內容結構（Anna 應該分幾個章節、每個章節講咩）\n"
            f"⑥ 語氣同風格要求\n\n"
            f"輸出一份清晰、完整嘅簡報，讓 Anna 可以完全根據呢份簡報製作最終成品，唔需要再猜或補充。"
        )
        _, structured_brief = await loop.run_in_executor(
            executor, run_with_system, AGENT_PROMPTS['Amy'], brief_prompt, AGENT_MODELS['Amy']
        )
        if "[QUOTA_EXCEEDED]" not in structured_brief:
            combined_context = f"【Amy 製作簡報】\n{structured_brief}"

    for d in file_dispatches:
        file_result = await maybe_generate_file(update, d["agent"], d["task"], loop, extra_context=combined_context)
        if file_result:
            agent_name, reply_text = file_result
            full_results += f"{AGENT_EMOJI.get(agent_name,'🤖')} {agent_name}：\n{reply_text}\n\n"
        else:
            enriched_task = (f"{combined_context}\n\n原始任務：\n{d['task']}") if combined_context else d["task"]
            _, reply_text = await loop.run_in_executor(executor, agent_call, d["agent"], enriched_task)
            emoji = AGENT_EMOJI.get(d["agent"], "🤖")
            msg = f"{emoji} {d['agent']}：\n{reply_text}"
            full_results += msg + "\n\n"
            await send_long(update, msg)

    all_results = (action_context + full_results).strip()
    if all_results:
        consolidate_user = (
            f"Stanley 嘅原始指令：{user_message}\n\n"
            f"各員工完整成果：\n{all_results[:5000]}\n\n"
            f"你係 Amy，Stanley 嘅秘書。所有員工已完成工作，請用廣東話出最終行動報告：\n"
            f"① 今次最重要嘅發現或成果（1-2句）\n"
            f"② 最優先行動（本週可以即時做嘅3件事，每件要具體）\n"
            f"③ 下一步建議（如需跟進或製作成品，說明點做）\n"
            f"格式清晰，唔超過200字。"
        )
        summary = await loop.run_in_executor(executor, run_with_system, AGENT_PROMPTS['Amy'], consolidate_user, AGENT_MODELS['Amy'])
        if "[QUOTA_EXCEEDED]" in summary:
            await update.message.reply_text(QUOTA_EXCEEDED_MSG)
            return
        summary_msg = f"{AGENT_EMOJI['Amy']} Amy 最終報告：\n\n{summary}"
        await send_long(update, summary_msg)
    else:
        summary_msg = ""

    # Save all agent outputs to last_content so follow-up "optimize/redo" commands have full context
    if all_results:
        last_content[ALLOWED_USER_ID] = all_results[:8000]
        save_last_content_to_disk(last_content)

    full_reply = f"{AGENT_EMOJI['Amy']} Amy：{amy_msg}\n\n{full_results}\n{summary_msg}".strip()
    history.append({"role": "user", "content": user_message})
    history.append({"role": "assistant", "content": full_reply})
    save_history(conversation_history)

    warning = get_usage_warning()
    if warning:
        await update.message.reply_text(warning)


# ── Daily report (with real web search) ───────────────────────────────────────

async def send_daily_report(app):
    try:
        from datetime import datetime as _dt
        loop = asyncio.get_event_loop()
        date_str = _dt.now().strftime("%Y年%m月%d日")

        import asyncio as _aio
        # RSS 係免費實時主力；DuckDuckGo 做 backup
        hk_news_rss, ai_raw_rss, industry_rss = await asyncio.gather(
            loop.run_in_executor(executor, rss_hk_news),
            loop.run_in_executor(executor, rss_ai_news),
            loop.run_in_executor(executor, rss_industry_news),
        )
        hk_news = hk_news_rss or ""
        ai_raw = ai_raw_rss or ""
        industry_raw = industry_rss or ""
        search_ok = WEB_SEARCH_AVAILABLE or (APIFY_AVAILABLE and APIFY_TOKEN)
        biz_news = ""
        if search_ok:
            if not hk_news:
                hk_news = await loop.run_in_executor(executor, news_search, "香港今日新聞 頭條 最新", 5)
                await _aio.sleep(1)
            biz_news = await loop.run_in_executor(executor, news_search, "香港財經商業 重要消息", 4)
            await _aio.sleep(1)
            if not ai_raw:
                ai_raw = await loop.run_in_executor(executor, news_search, "Claude ChatGPT Gemini AI 最新功能 更新", 5)
                await _aio.sleep(1)
            if not industry_raw:
                industry_raw = await loop.run_in_executor(executor, news_search, "香港美容 痛症治療 行業動態", 3)

        has_real_data = bool(hk_news or biz_news or ai_raw or industry_raw)
        if has_real_data:
            news_task = (
                f"根據以下今日香港真實新聞搜尋結果，列出5條最重要頭條，每條一句摘要：\n\n"
                f"{hk_news[:2000]}\n\n{biz_news[:1500]}\n\n"
                f"格式（每條一行）：📌 [標題]：[一句摘要]"
            )
            leo_task = (
                f"根據以下搜尋結果，生成香港美容同痛症業務情報（150字內）：\n\n"
                f"{industry_raw[:2000]}\n\n"
                f"輸出：① 今日最重要行業動態 ② 機會或風險 ③ Stanley 今日最應行動嘅一件事"
            )
            kai_task = (
                f"根據以下今日真實 AI 資訊，生成 AI 簡報（150字內）：\n\n"
                f"{ai_raw[:2000]}\n\n"
                f"輸出：① 最重要 AI 動態（只講事實）"
                f"② 對 Stanley 美容/痛症業務嘅直接影響 "
                f"③ 如有免費可用工具：「工具名 — 免費 — 可用於XX」；付費工具略過唔提"
            )
            search_tag = "📡 RSS + 🔍 實時搜尋"
        else:
            news_task = (
                f"今日（{date_str}）香港5條重要新聞頭條，每條一句摘要。"
                f"格式：📌 [標題]：[一句摘要]"
            )
            leo_task = (
                "今日香港美容同痛症業務情報（150字內）：\n"
                "① 最新行業趨勢或熱話 ② 競品動態 ③ Stanley 今日最應行動嘅一件事"
            )
            kai_task = (
                "今日全球 AI 市場簡報（150字內）：\n"
                "① 最新 AI 工具或功能更新（只講事實）\n"
                "② AI 用於美容/健康/銷售嘅直接影響\n"
                "③ 如有免費工具：「工具名 — 免費 — 可用於XX」；付費工具略過唔提"
            )
            search_tag = "📚 知識庫"

        agent_tasks = [
            loop.run_in_executor(executor, agent_call, "Leo", news_task),
            loop.run_in_executor(executor, agent_call, "Leo", leo_task),
            loop.run_in_executor(executor, agent_call, "Kai", kai_task),
        ]
        (_, headlines), (_, market), (_, ai_report) = await asyncio.gather(*agent_tasks)

        msg = (
            f"🌅 Stanley，早晨！{date_str} 情報簡報（{search_tag}）\n"
            f"━━━━━━━━━━━━━━\n\n"
            f"📰 今日頭條新聞：\n{headlines}\n\n"
            f"📊 Leo — 行業要聞：\n{market}\n\n"
            f"🤖 Kai — AI 資訊：\n{ai_report}\n\n"
            f"━━━━━━━━━━━━━━\n有咩問題隨時問我！"
        )

        for i in range(0, len(msg), 4096):
            await app.bot.send_message(chat_id=ALLOWED_USER_ID, text=msg[i:i+4096])

    except Exception as e:
        logger.error(f"Daily report error: {e}")


# ── /addcompetitor ─────────────────────────────────────────────────────────────

async def cmd_addcompetitor(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    args = context.args
    # Usage: /addcompetitor ig beautysignaturehk
    #        /addcompetitor maps "Beauty Signature HK"
    #        /addcompetitor xhs 香港美容
    if len(args) < 2:
        current = competitors_db
        lines = ["📋 **當前競品名單**\n"]
        for k, v in current.items():
            emoji = {"ig": "📸", "maps": "📍", "xhs": "🇨🇳", "threads": "🧵", "fb": "📘"}.get(k, "•")
            lines.append(f"{emoji} **{k.upper()}**: {', '.join(v) if v else '（空）'}")
        lines.append("\n用法：`/addcompetitor ig 帳號名`")
        lines.append("平台：ig / maps / xhs / threads / fb")
        await update.message.reply_text("\n".join(lines))
        return

    platform = args[0].lower()
    name = " ".join(args[1:]).strip().lstrip("@")

    if platform not in competitors_db:
        await update.message.reply_text(f"❌ 未知平台：{platform}\n可用：ig / maps / xhs / threads / fb")
        return

    if name not in competitors_db[platform]:
        competitors_db[platform].append(name)
        save_competitors(competitors_db)
        emoji = {"ig": "📸", "maps": "📍", "xhs": "🇨🇳", "threads": "🧵", "fb": "📘"}.get(platform, "•")
        await update.message.reply_text(f"✅ 已加入 {emoji} {platform.upper()} 名單：**{name}**\n\n用 `/intel` 做完整競品分析。")
    else:
        await update.message.reply_text(f"ℹ️ {name} 已喺 {platform.upper()} 名單入面。")


# ── /intel ─────────────────────────────────────────────────────────────────────

async def cmd_intel(update: Update, context: ContextTypes.DEFAULT_TYPE):
    if update.effective_user.id != ALLOWED_USER_ID:
        return

    comps = competitors_db
    ig_list = comps.get("ig", [])
    maps_list = comps.get("maps", [])
    xhs_list = comps.get("xhs", ["香港美容療程", "香港痛症"])

    if not ig_list and not maps_list:
        await update.message.reply_text(
            "⚠️ 競品名單係空嘅！先加入競品：\n"
            "`/addcompetitor ig beautysignaturehk`\n"
            "`/addcompetitor maps 美容診所 香港`\n\n"
            "加完再 `/intel` 即可。"
        )
        return

    await update.message.reply_text(
        f"🕵️ **競品情報掃描開始**\n\n"
        f"📸 IG 競品：{', '.join(ig_list) if ig_list else '（空）'}\n"
        f"📍 Maps 競品：{', '.join(maps_list) if maps_list else '（空）'}\n"
        f"🇨🇳 小紅書關鍵詞：{', '.join(xhs_list)}\n\n"
        f"預計需時 2-4 分鐘，完成後一次過輸出分析..."
    )

    loop = asyncio.get_event_loop()
    collected: dict[str, str] = {}

    # ── 並行抓取所有資料來源 ────────────────────────────────────────
    tasks = {}

    async def run_scrape(key: str, fn, *args):
        try:
            result = await loop.run_in_executor(executor, fn, *args)
            collected[key] = result or "（無數據）"
        except Exception as e:
            collected[key] = f"（抓取失敗：{e}）"

    scrape_tasks = []
    for ig in ig_list:
        scrape_tasks.append(run_scrape(f"ig_{ig}", scrape_instagram, ig, 15))
    for maps_q in maps_list:
        scrape_tasks.append(run_scrape(f"maps_{maps_q}", scrape_google_maps_reviews, maps_q, 25))
    for xhs_q in xhs_list:
        scrape_tasks.append(run_scrape(f"xhs_{xhs_q}", scrape_xiaohongshu, xhs_q, 10))

    await asyncio.gather(*scrape_tasks)

    # ── 過濾失敗嘅抓取，唔把廢數據餵俾員工 ──────────────────────────
    _fail_markers = ["⚠️", "FAILED", "抓取失敗", "（無數據）"]
    failed_keys = [k for k, v in collected.items() if any(m in v for m in _fail_markers)]
    ok_keys = [k for k in collected if k not in failed_keys]

    if failed_keys:
        fail_names = []
        for k in failed_keys:
            if k.startswith("ig_"):
                fail_names.append(f"IG @{k[3:]}")
            elif k.startswith("maps_"):
                fail_names.append(f"Maps: {k[5:]}")
            elif k.startswith("xhs_"):
                fail_names.append(f"小紅書: {k[4:]}")
            else:
                fail_names.append(k)
        await update.message.reply_text(
            "⚠️ 以下來源抓取失敗，已略過（節省 tokens）：\n" +
            "\n".join(f"• {n}" for n in fail_names) +
            "\n\n如係 IG/小紅書，請確認 Apify credits 是否充足。"
        )

    if not ok_keys:
        await update.message.reply_text("❌ 所有數據來源均抓取失敗，分析取消。\n請確認 Apify credits 充足後再試。")
        return

    # ── 整合成功嘅原始數據 ──────────────────────────────────────────
    raw_data = ""
    if any(k.startswith("ig_") for k in ok_keys):
        raw_data += "【競品 IG 帖文數據】\n"
        for k in ok_keys:
            if k.startswith("ig_"):
                raw_data += f"\n--- @{k[3:]} ---\n{collected[k][:1500]}\n"

    if any(k.startswith("maps_") for k in ok_keys):
        raw_data += "\n\n【Google Maps 評論（競品差評 = 市場空白位）】\n"
        for k in ok_keys:
            if k.startswith("maps_"):
                raw_data += f"\n--- {k[5:]} ---\n{collected[k][:1500]}\n"

    if any(k.startswith("xhs_") for k in ok_keys):
        raw_data += "\n\n【小紅書趨勢內容】\n"
        for k in ok_keys:
            if k.startswith("xhs_"):
                raw_data += f"\n--- 關鍵詞：{k[4:]} ---\n{collected[k][:1000]}\n"

    # ── Leo 做市場空白位分析 ────────────────────────────────────────
    await update.message.reply_text("📊 數據收集完成！Leo 分析中...")

    leo_prompt = f"""你係 Leo，市場情報分析師。根據以下真實競品數據，做深度市場空白位分析（Market Gap Analysis）。

{raw_data[:6000]}

請輸出：

## 1. 競品內容策略解碼
（每個競品：最常做咩類型內容、爆款規律、話術風格）

## 2. 市場空白位（最重要）
（競品完全唔做或做得差嘅主題/角度——Stanley 嘅機會）

## 3. 客戶最大痛點（來自差評）
（Google Maps 差評裡面，客人最不滿嘅係咩——Stanley 可以針對呢些痛點建立優勢）

## 4. 小紅書趨勢 vs Stanley 現況
（什麼內容正在爆紅但競品未做）

## 5. 本週立即可行動 3 件事
（具體到可以今天開始執行）

用廣東話輸出，有數據支撐，直接講結論。"""

    try:
        leo_result = await loop.run_in_executor(
            executor,
            lambda: anthropic.Anthropic(api_key=ANTHROPIC_API_KEY).messages.create(
                model="claude-sonnet-4-6",
                max_tokens=2000,
                system=AGENT_PROMPTS["Leo"],
                messages=[{"role": "user", "content": leo_prompt}]
            ).content[0].text
        )
    except Exception as e:
        leo_result = f"Leo 分析失敗：{e}"

    # ── Small 補充策略建議 ──────────────────────────────────────────
    small_prompt = f"""根據以下 Leo 嘅市場空白位分析，俾 Stanley 3個最優先嘅策略行動建議。
每個建議要有：① 做咩 ② 點做（具體步驟）③ 預期效果 ④ 風險提示

Leo 分析：
{leo_result[:2000]}

廣東話，挑戰者模式，唔好廢話。"""

    try:
        small_result = await loop.run_in_executor(
            executor,
            lambda: anthropic.Anthropic(api_key=ANTHROPIC_API_KEY).messages.create(
                model="claude-sonnet-4-6",
                max_tokens=800,
                system=AGENT_PROMPTS["Small"],
                messages=[{"role": "user", "content": small_prompt}]
            ).content[0].text
        )
    except Exception as e:
        small_result = f"Small 分析失敗：{e}"

    # ── 輸出完整報告 ────────────────────────────────────────────────
    ig_count = len(ig_list)
    maps_count = len(maps_list)
    await update.message.reply_text(
        f"🕵️ *競品情報報告*\n"
        f"抓取來源：{ig_count}個IG + {maps_count}個Maps + {len(xhs_list)}個小紅書關鍵詞",
        parse_mode="Markdown"
    )
    await send_with_avatar(update, "Leo", leo_result)
    await send_with_avatar(update, "Small", small_result)


# ── 客戶智識庫指令 ────────────────────────────────────────────────────────────────

async def cmd_add(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """
    /add [類別] [內容]
    類別：quote / hook / objection / win / gap
    或中文：客戶原話 / 開場白 / 反對意見 / 成功案例 / 空白位
    """
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    args = context.args or []
    if len(args) < 2:
        await update.message.reply_text(
            "📚 *客戶智識庫* — 存入格式：\n\n"
            "`/add quote 客戶話：腰痛咗三年，試過好多嘢都冇用`\n"
            "`/add hook 你係咪瞓醒就腰痛？`\n"
            "`/add objection 太貴 → 平均每次只係$X，少過一次按摩`\n"
            "`/add win 陳先生：2次療程後腰痛減輕7成`\n"
            "`/add gap 競品完全唔做「上班族久坐痛」角度`\n\n"
            "存入後，Anna/Tony/Rex 做嘢時自動使用。",
            parse_mode="Markdown"
        )
        return

    raw_cat = args[0]
    content = " ".join(args[1:]).strip()
    category = _INTEL_CATEGORY_ALIASES.get(raw_cat)
    if not category:
        await update.message.reply_text(
            f"❌ 唔識類別「{raw_cat}」\n\n"
            "可用類別：quote / hook / objection / win / gap\n"
            "中文：客戶原話 / 開場白 / 反對意見 / 成功案例 / 空白位"
        )
        return

    brand_intel_db[category].append(content)
    save_brand_intel(brand_intel_db)
    label = _INTEL_CATEGORY_LABELS[category]
    total = len(brand_intel_db[category])
    await update.message.reply_text(
        f"✅ 已存入 {label}\n\n"
        f"「{content[:100]}{'...' if len(content) > 100 else ''}」\n\n"
        f"該類別現有 {total} 條。Anna/Tony/Rex 下次做嘢時自動用到。"
    )


async def cmd_mydata(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """/mydata — 睇智識庫所有儲存內容"""
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    db = brand_intel_db
    total = sum(len(v) for v in db.values())
    if total == 0:
        await update.message.reply_text(
            "📚 智識庫係空嘅。\n\n"
            "用 `/add quote 客戶原話` 開始存入資料。"
        )
        return

    lines = [f"📚 *客戶智識庫* — 共 {total} 條\n"]
    for cat, label in _INTEL_CATEGORY_LABELS.items():
        items = db.get(cat, [])
        if items:
            lines.append(f"\n{label}（{len(items)} 條）")
            for i, item in enumerate(items, 1):
                preview = item[:80] + ("..." if len(item) > 80 else "")
                lines.append(f"  {i}. {preview}")
    lines.append("\n\n刪除某條：`/removedata [類別] [編號]`")
    await update.message.reply_text("\n".join(lines), parse_mode="Markdown")


async def cmd_removedata(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """/removedata [類別] [編號] — 刪除智識庫某條"""
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    args = context.args or []
    if len(args) < 2:
        await update.message.reply_text("用法：`/removedata quote 2` （刪除 quote 第2條）", parse_mode="Markdown")
        return
    raw_cat = args[0]
    category = _INTEL_CATEGORY_ALIASES.get(raw_cat)
    if not category:
        await update.message.reply_text(f"❌ 唔識類別「{raw_cat}」")
        return
    try:
        idx = int(args[1]) - 1
        items = brand_intel_db[category]
        if idx < 0 or idx >= len(items):
            await update.message.reply_text(f"❌ 編號超出範圍，該類別共 {len(items)} 條")
            return
        removed = items.pop(idx)
        save_brand_intel(brand_intel_db)
        label = _INTEL_CATEGORY_LABELS[category]
        await update.message.reply_text(f"🗑 已刪除 {label} 第 {idx+1} 條：\n「{removed[:100]}」")
    except ValueError:
        await update.message.reply_text("❌ 編號要係數字，例如：`/removedata quote 2`", parse_mode="Markdown")


# ── 全平台掃描 /scan ──────────────────────────────────────────────────────────────

async def cmd_scan(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """/scan [關鍵詞] — 同步抓取所有平台 + 5個分析角度，Leo 做跨平台整合分析"""
    if update.effective_user.id != ALLOWED_USER_ID:
        return
    if not context.args:
        await update.message.reply_text(
            "用法：`/scan [關鍵詞]`\n\n"
            "例子：\n"
            "`/scan 香港痛症管理`\n"
            "`/scan 美容療程`\n\n"
            "自動抓取：Google新聞 + 小紅書 + Threads + Facebook\n"
            "自動分析角度：客戶痛點 / 競品話術 / 真實評價 / 市場趨勢 / 客戶疑慮\n"
            "結果直接可用於 PDF / Landing Page / PPT 製作。",
            parse_mode="Markdown"
        )
        return

    keyword = " ".join(context.args).strip()
    loop = asyncio.get_event_loop()

    await update.message.reply_text(
        f"🔍 全平台掃描「{keyword}」\n\n"
        f"平台：Google新聞 / 小紅書 / Threads / Facebook\n"
        f"角度：客戶痛點 / 競品話術 / 真實評價 / 趨勢 / 疑慮\n\n"
        f"預計 45–90 秒，完成後 Leo 整合分析..."
    )

    _fail_kw = ["⚠️", "FAILED", "抓取失敗", "失敗", "Error", "error", "（無數據）", "未設定"]

    async def _fetch(name: str, fn, *args):
        try:
            result = await loop.run_in_executor(executor, fn, *args)
            failed = any(k in (result or "") for k in _fail_kw) or not (result or "").strip()
            return name, result or "", failed
        except Exception as e:
            return name, f"抓取失敗：{e}", True

    # ── 並行抓取：平台 + 分析角度 ──────────────────────────────────
    all_results = await asyncio.gather(
        # 平台數據
        _fetch("Google 新聞", scrape_google, f"{keyword} 香港 最新"),
        _fetch("小紅書",      scrape_xhs,    keyword),
        _fetch("Threads",    web_search,    f"site:threads.net {keyword}", 5),
        _fetch("Facebook",   web_search,    f"site:facebook.com {keyword} 香港", 5),
        # 分析角度（DuckDuckGo，免費穩定）
        _fetch("客戶痛點",   web_search, f"{keyword} 問題 點解 香港 煩惱", 5),
        _fetch("競品話術",   web_search, f"{keyword} 推廣 療程 香港 服務", 5),
        _fetch("真實評價",   web_search, f"{keyword} 效果 評價 真實 用家", 5),
        _fetch("2026趨勢",   web_search, f"{keyword} 2026 趨勢 香港 市場", 4),
        _fetch("客戶疑慮",   web_search, f"{keyword} 貴 值唔值 效果 安全", 4),
    )

    # ── 分類狀態報告 ───────────────────────────────────────────────
    platform_names = {"Google 新聞", "小紅書", "Threads", "Facebook"}
    platform_lines, angle_lines = [], []
    raw_platform, raw_angles = [], []

    for name, result, failed in all_results:
        tick = "✅" if not failed else "❌"
        if name in platform_names:
            platform_lines.append(f"{tick} {name}")
            if not failed:
                raw_platform.append(f"【{name}】\n{result[:800]}")
        else:
            angle_lines.append(f"{tick} {name}")
            if not failed:
                raw_angles.append(f"【{name}】\n{result[:600]}")

    await update.message.reply_text(
        f"📊 抓取完成：\n"
        f"平台：{'  '.join(platform_lines)}\n"
        f"角度：{'  '.join(angle_lines)}"
    )

    if not raw_platform and not raw_angles:
        await update.message.reply_text("❌ 所有來源均失敗，無法分析。請稍後再試。")
        return

    # ── Leo 跨平台整合分析 ─────────────────────────────────────────
    combined = "\n\n".join(raw_platform + raw_angles)
    leo_prompt = f"""以下係「{keyword}」嘅跨平台真實數據，請做深度整合分析。分析結果將用於製作 PDF、Landing Page 同 PPT，所以每個部分都要有可以直接引用嘅素材。

{combined[:6000]}

請輸出：

## 1. 客戶核心痛點（最重要）
（直接引用或還原客戶原話，用佢哋自己嘅字眼——呢啲係文稿嘅 hook 素材）

## 2. 跨平台共同趨勢
（多個平台同時出現嘅話題 = 真實市場信號）

## 3. 競品話術解碼
（競品點講、用咩字眼——Stanley 點差異化）

## 4. 社交證明素材
（真實評價角度、客戶最常分享咩結果——Landing Page 用）

## 5. 客戶最大疑慮 + 化解方向
（佢哋最怕咩、最常問咩——PDF FAQ 同 Landing Page objection section 用）

## 6. 市場空白位
（有人問但競品唔做嘅角度——Stanley 可以搶佔）

## 7. 本週立即可行動 3 件事

廣東話，有數據有原話，直接可以剪下貼入成品。"""

    await update.message.reply_text("🧠 Leo 整合分析緊，稍等...")
    _, leo_result = await loop.run_in_executor(executor, agent_call, "Leo", leo_prompt)
    await send_with_avatar(update, "Leo", leo_result)


# ── Main ───────────────────────────────────────────────────────────────────────

async def post_init(app):
    scheduler = AsyncIOScheduler(timezone="Asia/Hong_Kong")
    scheduler.add_job(send_daily_report, CronTrigger(hour=8, minute=45), args=[app])
    scheduler.start()


def main():
    app = Application.builder().token(BOT_TOKEN).post_init(post_init).build()

    app.add_handler(CommandHandler("start", cmd_start))
    app.add_handler(CommandHandler("clear", cmd_clear))
    app.add_handler(CommandHandler("status", cmd_status))
    app.add_handler(CommandHandler("export", cmd_export))
    app.add_handler(CommandHandler("dreamteam", cmd_dreamteam))
    app.add_handler(CommandHandler("dreamteampdf", cmd_dreamteampdf))
    app.add_handler(CommandHandler("reportpdf", cmd_reportpdf))
    app.add_handler(CommandHandler("landingpage", cmd_landingpage))
    app.add_handler(CommandHandler("pdf", cmd_pdf))
    app.add_handler(CommandHandler("slides", cmd_slides))
    app.add_handler(CommandHandler("caption", cmd_caption))
    app.add_handler(CommandHandler("reel", cmd_reel))
    app.add_handler(CommandHandler("copy", cmd_copy))
    app.add_handler(CommandHandler("report", cmd_report))
    app.add_handler(CommandHandler("imagine", cmd_imagine))
    app.add_handler(CommandHandler("scrape", cmd_scrape))
    app.add_handler(CommandHandler("xhs", cmd_xhs))
    app.add_handler(CommandHandler("research", cmd_research))
    app.add_handler(CommandHandler("intel", cmd_intel))
    app.add_handler(CommandHandler("addcompetitor", cmd_addcompetitor))
    app.add_handler(CommandHandler("done", cmd_done))
    app.add_handler(CommandHandler("myreport", cmd_myreport))
    app.add_handler(CommandHandler("add", cmd_add))
    app.add_handler(CommandHandler("mydata", cmd_mydata))
    app.add_handler(CommandHandler("removedata", cmd_removedata))
    app.add_handler(CommandHandler("scan", cmd_scan))
    app.add_handler(CallbackQueryHandler(handle_followup_callback, pattern="^followup:"))
    app.add_handler(CallbackQueryHandler(handle_report_action_callback, pattern="^report_action:"))
    app.add_handler(MessageHandler(filters.PHOTO, handle_photo))
    app.add_handler(MessageHandler(filters.Document.ALL, handle_document))
    app.add_handler(MessageHandler(filters.VOICE | filters.AUDIO, handle_voice))
    app.add_handler(MessageHandler(filters.TEXT & ~filters.COMMAND, handle_message))

    logger.info("Multi-Agent Bot v4.0 啟動")
    app.run_polling(allowed_updates=Update.ALL_TYPES)


if __name__ == "__main__":
    main()
